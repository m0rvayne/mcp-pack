"""
Google Workspace MCP — shared utilities.
Error formatting, file helpers, Sheets helpers, content readers.
"""

import io
import re
import base64
import zipfile
from pathlib import Path
from typing import Optional

from googleapiclient.http import MediaIoBaseDownload

BASE_DIR = Path(__file__).parent
DOWNLOAD_DIR = BASE_DIR / "downloads"
MAX_DOWNLOAD_SIZE = 100 * 1024 * 1024  # 100 MB

MIME_LABELS = {
    "application/vnd.google-apps.document": "Google Doc",
    "application/vnd.google-apps.spreadsheet": "Google Sheet",
    "application/vnd.google-apps.presentation": "Google Slides",
    "application/vnd.google-apps.form": "Google Form",
    "application/vnd.google-apps.folder": "Folder",
    "application/pdf": "PDF",
    "text/plain": "Text",
    "text/csv": "CSV",
    "text/html": "HTML",
    "text/markdown": "Markdown",
    "application/json": "JSON",
    "image/jpeg": "JPEG",
    "image/png": "PNG",
    "image/gif": "GIF",
    "image/webp": "WebP",
    "image/svg+xml": "SVG",
    "video/mp4": "MP4",
    "video/quicktime": "MOV",
    "audio/mpeg": "MP3",
    "application/zip": "ZIP",
    "application/x-zip-compressed": "ZIP",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "Excel (.xlsx)",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "Word (.docx)",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "PowerPoint (.pptx)",
    "application/vnd.ms-excel": "Excel (.xls)",
    "application/msword": "Word (.doc)",
}

EXPORT_FORMATS = {
    "application/vnd.google-apps.document": (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document", ".docx"),
    "application/vnd.google-apps.spreadsheet": (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", ".xlsx"),
    "application/vnd.google-apps.presentation": (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation", ".pptx"),
    "application/vnd.google-apps.drawing": ("image/png", ".png"),
}


# =============================================================================
# Error helpers
# =============================================================================

def safe_error(e: Exception) -> str:
    msg = str(e)
    msg = re.sub(r'/Users/[^\s:\"\']+', '<path>', msg)
    msg = re.sub(r'(key|token|password|secret|cookie)=[^\s&\"\']+', r'\1=<redacted>', msg, flags=re.IGNORECASE)
    msg = re.sub(r'Bearer\s+[^\s\"\']+', 'Bearer <redacted>', msg)
    return msg


def format_error(error: Exception) -> str:
    from googleapiclient.errors import HttpError
    if isinstance(error, HttpError):
        status = error.resp.status if hasattr(error, 'resp') else None
        reason = error._get_reason() if hasattr(error, '_get_reason') else str(error)
        if status == 403: return f"Access denied (HTTP 403): {reason}"
        if status == 404: return f"Not found (HTTP 404): {reason}"
        if status == 429: return "Rate limit exceeded (HTTP 429). Please wait and try again."
        if status == 400: return f"Invalid request (HTTP 400): {reason}"
        return f"Google API error (HTTP {status}): {reason}"
    return safe_error(error)


# =============================================================================
# Drive file helpers
# =============================================================================

def file_link(file_id: str, mime: str) -> str:
    if mime == "application/vnd.google-apps.document":
        return f"https://docs.google.com/document/d/{file_id}/edit"
    if mime == "application/vnd.google-apps.spreadsheet":
        return f"https://docs.google.com/spreadsheets/d/{file_id}/edit"
    if mime == "application/vnd.google-apps.presentation":
        return f"https://docs.google.com/presentation/d/{file_id}/edit"
    return f"https://drive.google.com/file/d/{file_id}/view"


def fmt_size(size) -> str:
    if not size: return "—"
    b = int(size)
    for unit in ("B", "KB", "MB", "GB"):
        if b < 1024: return f"{b:.1f} {unit}"
        b /= 1024
    return f"{b:.1f} TB"


def fmt_file(f: dict) -> dict:
    mime = f.get("mimeType", "")
    return {
        "id": f["id"],
        "name": f.get("name", ""),
        "type": MIME_LABELS.get(mime, mime),
        "mimeType": mime,
        "size": fmt_size(f.get("size")),
        "modified": f.get("modifiedTime", "")[:10],
        "link": file_link(f["id"], mime),
    }


def download_bytes(service, file_id: str, mime: str) -> bytes:
    meta = service.files().get(fileId=file_id, fields="size", supportsAllDrives=True).execute()
    size = int(meta.get("size") or 0)
    if size > MAX_DOWNLOAD_SIZE:
        raise ValueError(f"File too large ({fmt_size(size)}). Max: {fmt_size(MAX_DOWNLOAD_SIZE)}")
    request = service.files().get_media(fileId=file_id)
    buf = io.BytesIO()
    dl = MediaIoBaseDownload(buf, request)
    done = False
    while not done:
        _, done = dl.next_chunk()
    return buf.getvalue()


def export_bytes(service, file_id: str, export_mime: str) -> bytes:
    request = service.files().export_media(fileId=file_id, mimeType=export_mime)
    buf = io.BytesIO()
    dl = MediaIoBaseDownload(buf, request)
    done = False
    while not done:
        _, done = dl.next_chunk()
    return buf.getvalue()


def truncate(text: str, limit: int = 10000) -> str:
    if len(text) <= limit: return text
    return text[:limit] + f"\n\n... [truncated — {len(text) - limit} more chars]"


# =============================================================================
# Sheets helpers
# =============================================================================

def col_letter_to_index(col: str) -> int:
    result = 0
    for ch in col.upper():
        result = result * 26 + (ord(ch) - ord("A") + 1)
    return result - 1


def col_index_to_letter(index: int) -> str:
    result = ""
    index += 1
    while index > 0:
        index, remainder = divmod(index - 1, 26)
        result = chr(65 + remainder) + result
    return result


def parse_a1_range(a1: str) -> dict:
    sheet = None
    range_part = a1
    if "!" in a1:
        sheet, range_part = a1.rsplit("!", 1)
        sheet = sheet.strip("'\"")
    match = re.match(r"^([A-Za-z]+)?(\d+)?(?::([A-Za-z]+)?(\d+)?)?$", range_part)
    if not match:
        return {"sheet": sheet, "raw": range_part}
    start_col, start_row, end_col, end_row = match.groups()
    return {
        "sheet": sheet,
        "start_col": start_col,
        "start_row": int(start_row) if start_row else None,
        "end_col": end_col,
        "end_row": int(end_row) if end_row else None,
    }


def validate_spreadsheet_id(spreadsheet_id: str) -> str:
    if "/" in spreadsheet_id:
        match = re.search(r"/spreadsheets/d/([a-zA-Z0-9-_]+)", spreadsheet_id)
        if match: return match.group(1)
        raise ValueError(f"Cannot extract spreadsheet ID from URL: {spreadsheet_id}")
    return spreadsheet_id.strip()


def make_color(hex_color: str) -> dict:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 6:
        r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
        return {"red": r / 255, "green": g / 255, "blue": b / 255}
    return {"red": 0, "green": 0, "blue": 0}


def make_text_format(**kwargs) -> dict:
    fmt = {}
    if "bold" in kwargs: fmt["bold"] = kwargs["bold"]
    if "italic" in kwargs: fmt["italic"] = kwargs["italic"]
    if "underline" in kwargs: fmt["underline"] = kwargs["underline"]
    if "strikethrough" in kwargs: fmt["strikethrough"] = kwargs["strikethrough"]
    if "font_size" in kwargs: fmt["fontSize"] = kwargs["font_size"]
    if "font_family" in kwargs: fmt["fontFamily"] = kwargs["font_family"]
    if "font_color" in kwargs: fmt["foregroundColorStyle"] = {"rgbColor": make_color(kwargs["font_color"])}
    return fmt


def make_cell_format(bg_color=None, h_align=None, v_align=None, wrap_strategy=None,
                     number_format=None, number_format_type=None, **text_kwargs):
    cell_format = {}
    fields = []
    if bg_color:
        cell_format["backgroundColor"] = make_color(bg_color)
        fields.append("userEnteredFormat.backgroundColor")
    if h_align:
        cell_format["horizontalAlignment"] = h_align.upper()
        fields.append("userEnteredFormat.horizontalAlignment")
    if v_align:
        cell_format["verticalAlignment"] = v_align.upper()
        fields.append("userEnteredFormat.verticalAlignment")
    if wrap_strategy:
        cell_format["wrapStrategy"] = wrap_strategy.upper()
        fields.append("userEnteredFormat.wrapStrategy")
    if number_format or number_format_type:
        nf = {}
        if number_format_type: nf["type"] = number_format_type.upper()
        if number_format: nf["pattern"] = number_format
        cell_format["numberFormat"] = nf
        fields.append("userEnteredFormat.numberFormat")
    text_fmt = make_text_format(**text_kwargs)
    if text_fmt:
        cell_format["textFormat"] = text_fmt
        fields.append("userEnteredFormat.textFormat")
    return cell_format, ",".join(fields)


def a1_to_grid_range(a1: str, sheet_id: int) -> dict:
    parsed = parse_a1_range(a1)
    gr = {"sheetId": sheet_id}
    if parsed.get("start_col"):
        gr["startColumnIndex"] = col_letter_to_index(parsed["start_col"])
    if parsed.get("end_col"):
        gr["endColumnIndex"] = col_letter_to_index(parsed["end_col"]) + 1
    if parsed.get("start_row"):
        gr["startRowIndex"] = parsed["start_row"] - 1
    if parsed.get("end_row"):
        gr["endRowIndex"] = parsed["end_row"]
    return gr


def rows_to_table(values, has_header=True):
    if not values:
        return {"headers": [], "rows": [], "total_rows": 0}
    if has_header:
        headers = values[0]
        data_rows = values[1:]
    else:
        max_cols = max(len(row) for row in values)
        headers = [col_index_to_letter(i) for i in range(max_cols)]
        data_rows = values
    rows = []
    for row in data_rows:
        row_dict = {}
        for i, header in enumerate(headers):
            row_dict[header] = row[i] if i < len(row) else ""
        rows.append(row_dict)
    return {"headers": headers, "rows": rows, "total_rows": len(rows)}


# =============================================================================
# Content readers
# =============================================================================

def read_pdf(data: bytes) -> str:
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(data))
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)
    except ImportError:
        return "[PDF reading requires PyPDF2: pip install PyPDF2]"


def read_docx(data: bytes) -> str:
    try:
        import docx
        doc = docx.Document(io.BytesIO(data))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join(parts)
    except ImportError:
        return "[DOCX reading requires python-docx: pip install python-docx]"


def read_xlsx(data: bytes) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Sheet: {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                parts.append(" | ".join(str(c) if c is not None else "" for c in row))
        return "\n".join(parts)
    except ImportError:
        return "[XLSX reading requires openpyxl: pip install openpyxl]"


def read_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
            if texts:
                parts.append(f"## Slide {i}\n" + "\n".join(texts))
        return "\n\n".join(parts)
    except ImportError:
        return "[PPTX reading requires python-pptx: pip install python-pptx]"


def read_zip(data: bytes) -> str:
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        files = []
        for info in zf.infolist():
            files.append(f"{'📁' if info.is_dir() else '📄'} {info.filename} ({fmt_size(info.file_size)})")
        return f"ZIP archive ({len(files)} items):\n" + "\n".join(files[:100])
