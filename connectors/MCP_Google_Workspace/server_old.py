"""
Google Workspace MCP Server
Unified server for Google Drive, Calendar, Sheets, Gmail, Docs, and Tasks.
Provides all tools through a single MCP endpoint with shared OAuth credentials.
"""

import asyncio
import io
import json
import base64
import os
import re as _re
import signal
import sys
import threading
import zipfile
from pathlib import Path
from typing import Any, Optional
from datetime import datetime, timedelta, timezone

from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp import types

from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload


# =============================================================================
# Signal handling
# =============================================================================

def _handle_shutdown(sig, frame):
    sys.exit(0)

signal.signal(signal.SIGTERM, _handle_shutdown)
signal.signal(signal.SIGINT, _handle_shutdown)


# =============================================================================
# Config
# =============================================================================

BASE_DIR = Path(__file__).parent
CREDENTIALS_FILE = BASE_DIR / "credentials.json"
TOKEN_FILE = BASE_DIR / "token.json"
DOWNLOAD_DIR = BASE_DIR / "downloads"

SCOPES = [
    "https://www.googleapis.com/auth/drive",
    "https://www.googleapis.com/auth/calendar",
    "https://www.googleapis.com/auth/spreadsheets",
    # TODO: Add these scopes and re-auth to enable Gmail, Docs, Tasks tools:
    # "https://www.googleapis.com/auth/gmail.modify",
    # "https://www.googleapis.com/auth/documents",
    # "https://www.googleapis.com/auth/tasks",
]

MIME_LABELS = {
    "application/vnd.google-apps.document": "Google Doc",
    "application/vnd.google-apps.spreadsheet": "Google Sheet",
    "application/vnd.google-apps.presentation": "Google Slides",
    "application/vnd.google-apps.form": "Google Form",
    "application/vnd.google-apps.folder": "Folder",
    "application/pdf": "PDF",
    "text/plain": "Text",
    "text/csv": "CSV",
    "text/html": "HTML",
    "text/markdown": "Markdown",
    "application/json": "JSON",
    "image/jpeg": "JPEG",
    "image/png": "PNG",
    "image/gif": "GIF",
    "image/webp": "WebP",
    "image/svg+xml": "SVG",
    "video/mp4": "MP4",
    "video/quicktime": "MOV",
    "audio/mpeg": "MP3",
    "application/zip": "ZIP",
    "application/x-zip-compressed": "ZIP",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "Excel (.xlsx)",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "Word (.docx)",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "PowerPoint (.pptx)",
    "application/vnd.ms-excel": "Excel (.xls)",
    "application/msword": "Word (.doc)",
}

MAX_DOWNLOAD_SIZE = 100 * 1024 * 1024  # 100 MB

EXPORT_FORMATS = {
    "application/vnd.google-apps.document": (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document", ".docx"),
    "application/vnd.google-apps.spreadsheet": (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", ".xlsx"),
    "application/vnd.google-apps.presentation": (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation", ".pptx"),
    "application/vnd.google-apps.drawing": ("image/png", ".png"),
}


# =============================================================================
# Shared Auth
# =============================================================================

_creds_cache = None
_creds_lock = threading.Lock()
_drive_service = None
_calendar_service = None
_sheets_service = None


def _get_creds():
    global _creds_cache, _drive_service, _calendar_service, _sheets_service
    with _creds_lock:
        if _creds_cache and _creds_cache.valid:
            return _creds_cache
        creds = None
        if TOKEN_FILE.exists():
            creds = Credentials.from_authorized_user_file(str(TOKEN_FILE), SCOPES)
        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                creds.refresh(Request())
            else:
                if not CREDENTIALS_FILE.exists():
                    raise RuntimeError(
                        "Google authentication required. Place credentials.json in "
                        f"{BASE_DIR} and run setup_auth.py first."
                    )
                flow = InstalledAppFlow.from_client_secrets_file(
                    str(CREDENTIALS_FILE), SCOPES
                )
                creds = flow.run_local_server(port=8085, prompt="consent", access_type="offline")
            TOKEN_FILE.write_text(creds.to_json())
            os.chmod(TOKEN_FILE, 0o600)
            # Reset cached services so they pick up new credentials
            _drive_service = None
            _calendar_service = None
            _sheets_service = None
        _creds_cache = creds
        return creds


def get_drive_service():
    global _drive_service
    if _drive_service is None:
        _drive_service = build("drive", "v3", credentials=_get_creds())
    return _drive_service


def get_calendar_service():
    global _calendar_service
    if _calendar_service is None:
        _calendar_service = build("calendar", "v3", credentials=_get_creds())
    return _calendar_service


def get_sheets_service():
    global _sheets_service
    if _sheets_service is None:
        _sheets_service = build("sheets", "v4", credentials=_get_creds())
    return _sheets_service


# =============================================================================
# Error helpers
# =============================================================================

def _safe_error(e: Exception) -> str:
    """Sanitize error messages to prevent leaking sensitive info."""
    msg = str(e)
    msg = _re.sub(r'/Users/[^\s:\"\']+', '<path>', msg)
    msg = _re.sub(r'(key|token|password|secret|cookie)=[^\s&\"\']+', r'\1=<redacted>', msg, flags=_re.IGNORECASE)
    msg = _re.sub(r'Bearer\s+[^\s\"\']+', 'Bearer <redacted>', msg)
    return msg


def _format_error(error: Exception) -> str:
    """Format API errors into human-readable messages."""
    if isinstance(error, HttpError):
        status = error.resp.status if hasattr(error, 'resp') else None
        reason = error._get_reason() if hasattr(error, '_get_reason') else str(error)
        if status == 403:
            return f"Access denied (HTTP 403): {reason}"
        if status == 404:
            return f"Not found (HTTP 404): {reason}"
        if status == 429:
            return "Rate limit exceeded (HTTP 429). Please wait and try again."
        if status == 400:
            return f"Invalid request (HTTP 400): {reason}"
        return f"Google API error (HTTP {status}): {reason}"
    return _safe_error(error)


# =============================================================================
# Sheets helpers (inlined from helpers.py)
# =============================================================================

def _col_letter_to_index(col: str) -> int:
    result = 0
    for ch in col.upper():
        result = result * 26 + (ord(ch) - ord("A") + 1)
    return result - 1


def _col_index_to_letter(index: int) -> str:
    result = ""
    index += 1
    while index > 0:
        index, remainder = divmod(index - 1, 26)
        result = chr(65 + remainder) + result
    return result


def _parse_a1_range(a1: str) -> dict:
    sheet = None
    range_part = a1
    if "!" in a1:
        sheet, range_part = a1.rsplit("!", 1)
        sheet = sheet.strip("'\"")
    match = _re.match(r"^([A-Za-z]+)?(\d+)?(?::([A-Za-z]+)?(\d+)?)?$", range_part)
    if not match:
        return {"sheet": sheet, "raw": range_part}
    start_col, start_row, end_col, end_row = match.groups()
    return {
        "sheet": sheet,
        "start_col": start_col,
        "start_row": int(start_row) if start_row else None,
        "end_col": end_col,
        "end_row": int(end_row) if end_row else None,
    }


def _validate_spreadsheet_id(spreadsheet_id: str) -> str:
    if "/" in spreadsheet_id:
        match = _re.search(r"/spreadsheets/d/([a-zA-Z0-9-_]+)", spreadsheet_id)
        if match:
            return match.group(1)
        raise ValueError(f"Cannot extract spreadsheet ID from URL: {spreadsheet_id}")
    return spreadsheet_id.strip()


def _make_color(hex_color: str) -> dict:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 6:
        r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
        return {"red": r / 255, "green": g / 255, "blue": b / 255}
    return {"red": 0, "green": 0, "blue": 0}


def _make_text_format(**kwargs) -> dict:
    fmt = {}
    if "bold" in kwargs: fmt["bold"] = kwargs["bold"]
    if "italic" in kwargs: fmt["italic"] = kwargs["italic"]
    if "underline" in kwargs: fmt["underline"] = kwargs["underline"]
    if "strikethrough" in kwargs: fmt["strikethrough"] = kwargs["strikethrough"]
    if "font_size" in kwargs: fmt["fontSize"] = kwargs["font_size"]
    if "font_family" in kwargs: fmt["fontFamily"] = kwargs["font_family"]
    if "font_color" in kwargs: fmt["foregroundColorStyle"] = {"rgbColor": _make_color(kwargs["font_color"])}
    return fmt


def _make_cell_format(
    bg_color=None, h_align=None, v_align=None, wrap_strategy=None,
    number_format=None, number_format_type=None, **text_kwargs,
):
    cell_format = {}
    fields = []
    if bg_color:
        cell_format["backgroundColor"] = _make_color(bg_color)
        fields.append("userEnteredFormat.backgroundColor")
    if h_align:
        cell_format["horizontalAlignment"] = h_align.upper()
        fields.append("userEnteredFormat.horizontalAlignment")
    if v_align:
        cell_format["verticalAlignment"] = v_align.upper()
        fields.append("userEnteredFormat.verticalAlignment")
    if wrap_strategy:
        cell_format["wrapStrategy"] = wrap_strategy.upper()
        fields.append("userEnteredFormat.wrapStrategy")
    if number_format or number_format_type:
        nf = {}
        if number_format_type: nf["type"] = number_format_type.upper()
        if number_format: nf["pattern"] = number_format
        cell_format["numberFormat"] = nf
        fields.append("userEnteredFormat.numberFormat")
    text_fmt = _make_text_format(**text_kwargs)
    if text_fmt:
        cell_format["textFormat"] = text_fmt
        fields.append("userEnteredFormat.textFormat")
    return cell_format, ",".join(fields)


def _a1_to_grid_range(a1: str, sheet_id: int) -> dict:
    parsed = _parse_a1_range(a1)
    gr = {"sheetId": sheet_id}
    if parsed.get("start_col"):
        gr["startColumnIndex"] = _col_letter_to_index(parsed["start_col"])
    if parsed.get("end_col"):
        gr["endColumnIndex"] = _col_letter_to_index(parsed["end_col"]) + 1
    if parsed.get("start_row"):
        gr["startRowIndex"] = parsed["start_row"] - 1
    if parsed.get("end_row"):
        gr["endRowIndex"] = parsed["end_row"]
    return gr


def _rows_to_table(values, has_header=True):
    if not values:
        return {"headers": [], "rows": [], "total_rows": 0}
    if has_header:
        headers = values[0]
        data_rows = values[1:]
    else:
        max_cols = max(len(row) for row in values)
        headers = [_col_index_to_letter(i) for i in range(max_cols)]
        data_rows = values
    rows = []
    for row in data_rows:
        row_dict = {}
        for i, header in enumerate(headers):
            row_dict[header] = row[i] if i < len(row) else ""
        rows.append(row_dict)
    return {"headers": headers, "rows": rows, "total_rows": len(rows)}


# =============================================================================
# Drive helpers
# =============================================================================

def file_link(file_id: str, mime: str) -> str:
    if mime == "application/vnd.google-apps.document":
        return f"https://docs.google.com/document/d/{file_id}/edit"
    if mime == "application/vnd.google-apps.spreadsheet":
        return f"https://docs.google.com/spreadsheets/d/{file_id}/edit"
    if mime == "application/vnd.google-apps.presentation":
        return f"https://docs.google.com/presentation/d/{file_id}/edit"
    return f"https://drive.google.com/file/d/{file_id}/view"


def fmt_size(size) -> str:
    if not size:
        return "—"
    b = int(size)
    for unit in ("B", "KB", "MB", "GB"):
        if b < 1024:
            return f"{b:.1f} {unit}"
        b /= 1024
    return f"{b:.1f} TB"


def fmt_file(f: dict) -> dict:
    mime = f.get("mimeType", "")
    return {
        "id": f["id"],
        "name": f.get("name", ""),
        "type": MIME_LABELS.get(mime, mime),
        "mimeType": mime,
        "size": fmt_size(f.get("size")),
        "modified": f.get("modifiedTime", "")[:10],
        "link": file_link(f["id"], mime),
    }


def download_bytes(service, file_id: str, mime: str) -> bytes:
    meta = service.files().get(fileId=file_id, fields="size", supportsAllDrives=True).execute()
    size = int(meta.get("size") or 0)
    if size > MAX_DOWNLOAD_SIZE:
        raise ValueError(f"File too large ({fmt_size(size)}). Max allowed: {fmt_size(MAX_DOWNLOAD_SIZE)}")
    request = service.files().get_media(fileId=file_id)
    buf = io.BytesIO()
    dl = MediaIoBaseDownload(buf, request)
    done = False
    while not done:
        _, done = dl.next_chunk()
    return buf.getvalue()


def export_bytes(service, file_id: str, export_mime: str) -> bytes:
    request = service.files().export_media(fileId=file_id, mimeType=export_mime)
    buf = io.BytesIO()
    dl = MediaIoBaseDownload(buf, request)
    done = False
    while not done:
        _, done = dl.next_chunk()
    return buf.getvalue()


def truncate(text: str, limit: int = 10000) -> str:
    if len(text) <= limit:
        return text
    return text[:limit] + f"\n\n... [truncated — {len(text) - limit} more chars]"


# -- Content readers --

def read_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"--- Page {i} ---\n{text}")
    return "\n\n".join(pages) if pages else "(No extractable text found in PDF)"


def read_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        parts.append("\n".join(rows))
    return "\n\n".join(parts) if parts else "(Empty document)"


def read_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"=== Sheet: {sheet_name} ===")
        rows = []
        for row in ws.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                rows.append(" | ".join(str(c) if c is not None else "" for c in row))
        if rows:
            parts.append("\n".join(rows))
        else:
            parts.append("(empty sheet)")
    return "\n\n".join(parts)


def read_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
    return "\n\n".join(slides) if slides else "(No text found in presentation)"


def read_zip(data: bytes) -> str:
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        names = zf.namelist()
        infos = zf.infolist()
        lines = [f"ZIP archive — {len(names)} file(s):\n"]
        for info in infos:
            size = fmt_size(info.file_size)
            lines.append(f"  {info.filename}  ({size})")
        return "\n".join(lines)


# =============================================================================
# Calendar helpers
# =============================================================================

def _fmt_event(e: dict) -> dict:
    start = e.get("start", {})
    end = e.get("end", {})
    result = {
        "id": e.get("id"),
        "summary": e.get("summary", "(no title)"),
        "start": start.get("dateTime", start.get("date", "")),
        "end": end.get("dateTime", end.get("date", "")),
        "location": e.get("location"),
        "status": e.get("status"),
        "htmlLink": e.get("htmlLink"),
        "attendees": [{"email": a["email"], "response": a.get("responseStatus", "")} for a in e.get("attendees", [])],
        "description": e.get("description"),
        "colorId": e.get("colorId"),
        "recurrence": e.get("recurrence"),
        "creator": e.get("creator", {}).get("email"),
        "organizer": e.get("organizer", {}).get("email"),
    }
    conf = e.get("conferenceData")
    if conf:
        for ep in conf.get("entryPoints", []):
            if ep.get("entryPointType") == "video":
                result["conferenceLink"] = ep.get("uri")
                break
        solution = conf.get("conferenceSolution", {})
        if solution:
            result["conferenceType"] = solution.get("name")
    hangout = e.get("hangoutLink")
    if hangout and "conferenceLink" not in result:
        result["conferenceLink"] = hangout
    return result


def _build_event_body(a: dict) -> dict:
    body: dict = {}
    if "summary" in a: body["summary"] = a["summary"]
    if "description" in a: body["description"] = a["description"]
    if "location" in a: body["location"] = a["location"]
    if "color_id" in a: body["colorId"] = a["color_id"]
    if "status" in a: body["status"] = a["status"]
    if "visibility" in a: body["visibility"] = a["visibility"]
    if "recurrence" in a: body["recurrence"] = [a["recurrence"]]

    tz = a.get("time_zone")
    if "start" in a:
        if "T" in a["start"]:
            body["start"] = {"dateTime": a["start"]}
            if tz: body["start"]["timeZone"] = tz
        else:
            body["start"] = {"date": a["start"]}
    if "end" in a:
        if "T" in a["end"]:
            body["end"] = {"dateTime": a["end"]}
            if tz: body["end"]["timeZone"] = tz
        else:
            body["end"] = {"date": a["end"]}

    if "attendees" in a:
        body["attendees"] = [{"email": e.strip()} for e in a["attendees"].split(",") if e.strip()]

    if "reminders_minutes" in a:
        mins = [int(m.strip()) for m in a["reminders_minutes"].split(",") if m.strip()]
        body["reminders"] = {
            "useDefault": False,
            "overrides": [{"method": "popup", "minutes": m} for m in mins],
        }

    if "zoom_link" in a and a["zoom_link"]:
        zoom_url = a["zoom_link"].strip()
        m = _re.search(r"/j/(\d+)", zoom_url)
        meeting_id = m.group(1) if m else ""
        body["conferenceData"] = {
            "conferenceSolution": {
                "key": {"type": "addOn"},
                "name": "Zoom Meeting",
                "iconUri": "https://lh3.googleusercontent.com/ugWKMQWEsDNegkGTPBjHpiMBeIbFx42mLPHEHBjkCNOr8TOQZfMjPb4GpiEL0Cpb2A",
            },
            "entryPoints": [
                {
                    "entryPointType": "video",
                    "uri": zoom_url,
                    "label": f"Zoom Meeting {meeting_id}".strip(),
                },
            ],
        }

    return body


# =============================================================================
# Sheets service functions (inlined from sheets_service.py)
# =============================================================================

def _sheets():
    return get_sheets_service().spreadsheets()


def _sheets_get_spreadsheet_info(spreadsheet_id: str) -> dict:
    sid = _validate_spreadsheet_id(spreadsheet_id)
    sp = _sheets().get(
        spreadsheetId=sid,
        fields="spreadsheetId,properties,sheets.properties,namedRanges",
    ).execute()
    sheets_info = []
    for s in sp.get("sheets", []):
        p = s["properties"]
        sheets_info.append({
            "sheetId": p["sheetId"],
            "title": p["title"],
            "index": p["index"],
            "rows": p.get("gridProperties", {}).get("rowCount"),
            "columns": p.get("gridProperties", {}).get("columnCount"),
            "tabColor": p.get("tabColorStyle", {}).get("rgbColor"),
            "hidden": p.get("hidden", False),
        })
    return {
        "spreadsheetId": sp["spreadsheetId"],
        "title": sp["properties"]["title"],
        "locale": sp["properties"].get("locale"),
        "timeZone": sp["properties"].get("timeZone"),
        "sheets": sheets_info,
        "namedRanges": sp.get("namedRanges", []),
    }


def _sheets_list_sheets(spreadsheet_id: str) -> list[dict]:
    info = _sheets_get_spreadsheet_info(spreadsheet_id)
    return info["sheets"]


def _sheets_add_sheet(spreadsheet_id, title, rows=1000, columns=26, tab_color=None):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    props = {"title": title, "gridProperties": {"rowCount": rows, "columnCount": columns}}
    if tab_color:
        props["tabColorStyle"] = {"rgbColor": _make_color(tab_color)}
    result = _sheets().batchUpdate(
        spreadsheetId=sid,
        body={"requests": [{"addSheet": {"properties": props}}]},
    ).execute()
    new_props = result["replies"][0]["addSheet"]["properties"]
    return {
        "sheetId": new_props["sheetId"], "title": new_props["title"],
        "rows": new_props["gridProperties"]["rowCount"],
        "columns": new_props["gridProperties"]["columnCount"],
    }


def _sheets_delete_sheet(spreadsheet_id, sheet_id):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"deleteSheet": {"sheetId": sheet_id}}]}).execute()
    return {"deleted": True, "sheetId": sheet_id}


def _sheets_rename_sheet(spreadsheet_id, sheet_id, new_title):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{
        "updateSheetProperties": {"properties": {"sheetId": sheet_id, "title": new_title}, "fields": "title"}
    }]}).execute()
    return {"sheetId": sheet_id, "newTitle": new_title}


def _sheets_copy_sheet(spreadsheet_id, sheet_id, destination_spreadsheet_id=None):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    dest = _validate_spreadsheet_id(destination_spreadsheet_id) if destination_spreadsheet_id else sid
    result = _sheets().sheets().copyTo(
        spreadsheetId=sid, sheetId=sheet_id, body={"destinationSpreadsheetId": dest},
    ).execute()
    return {"newSheetId": result["sheetId"], "title": result["title"], "destinationSpreadsheetId": dest}


def _sheets_reorder_sheets(spreadsheet_id, sheet_order):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    requests = [{"updateSheetProperties": {"properties": {"sheetId": item["sheetId"], "index": item["index"]}, "fields": "index"}} for item in sheet_order]
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": requests}).execute()
    return {"reordered": True, "order": sheet_order}


def _sheets_read_range(spreadsheet_id, range_a1, value_render="FORMATTED_VALUE"):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    result = _sheets().values().get(
        spreadsheetId=sid, range=range_a1, valueRenderOption=value_render, dateTimeRenderOption="FORMATTED_STRING",
    ).execute()
    values = result.get("values", [])
    return {"range": result.get("range"), "values": values, "rows": len(values), "cols": max((len(r) for r in values), default=0)}


def _sheets_batch_read(spreadsheet_id, ranges, value_render="FORMATTED_VALUE"):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    result = _sheets().values().batchGet(spreadsheetId=sid, ranges=ranges, valueRenderOption=value_render).execute()
    output = []
    for vr in result.get("valueRanges", []):
        values = vr.get("values", [])
        output.append({"range": vr.get("range"), "values": values, "rows": len(values)})
    return output


def _sheets_get_cell(spreadsheet_id, cell_a1):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    formatted = _sheets().values().get(spreadsheetId=sid, range=cell_a1, valueRenderOption="FORMATTED_VALUE").execute()
    formula = _sheets().values().get(spreadsheetId=sid, range=cell_a1, valueRenderOption="FORMULA").execute()
    fmt_val = ""
    formula_val = ""
    if formatted.get("values"):
        fmt_val = formatted["values"][0][0] if formatted["values"][0] else ""
    if formula.get("values"):
        formula_val = formula["values"][0][0] if formula["values"][0] else ""
    has_formula = isinstance(formula_val, str) and formula_val.startswith("=")
    return {"cell": cell_a1, "value": fmt_val, "formula": formula_val if has_formula else None}


def _sheets_search_values(spreadsheet_id, query, sheet_name=None):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    info = _sheets_get_spreadsheet_info(sid)
    max_rows = 50000
    sheets_to_search = info["sheets"]
    if sheet_name:
        sheets_to_search = [s for s in sheets_to_search if s["title"] == sheet_name]
    results = []
    for sheet in sheets_to_search:
        title = sheet["title"]
        row_count = min(sheet.get("rows") or max_rows, max_rows)
        try:
            data = _sheets().values().get(spreadsheetId=sid, range=f"'{title}'!1:{row_count}", valueRenderOption="FORMATTED_VALUE").execute()
        except Exception:
            continue
        for row_idx, row in enumerate(data.get("values", [])):
            for col_idx, val in enumerate(row):
                if query.lower() in str(val).lower():
                    results.append({"sheet": title, "cell": f"{_col_index_to_letter(col_idx)}{row_idx + 1}", "value": val, "row": row_idx + 1, "column": _col_index_to_letter(col_idx)})
    return results


def _sheets_get_sheet_as_table(spreadsheet_id, sheet_name=None, has_header=True, max_rows=0):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    if max_rows <= 0 or max_rows > 50000:
        max_rows = 50000
    range_a1 = f"'{sheet_name}'" if sheet_name else ""
    if not range_a1:
        info = _sheets_get_spreadsheet_info(sid)
        if info["sheets"]:
            range_a1 = f"'{info['sheets'][0]['title']}'"
    result = _sheets().values().get(spreadsheetId=sid, range=range_a1, valueRenderOption="FORMATTED_VALUE").execute()
    values = result.get("values", [])
    if has_header:
        values = values[:max_rows + 1]
    else:
        values = values[:max_rows]
    return _rows_to_table(values, has_header)


def _sheets_write_range(spreadsheet_id, range_a1, values, value_input="USER_ENTERED"):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    result = _sheets().values().update(spreadsheetId=sid, range=range_a1, valueInputOption=value_input, body={"values": values}).execute()
    return {"updatedRange": result.get("updatedRange"), "updatedRows": result.get("updatedRows"), "updatedColumns": result.get("updatedColumns"), "updatedCells": result.get("updatedCells")}


def _sheets_append_rows(spreadsheet_id, range_a1, values, value_input="USER_ENTERED"):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    result = _sheets().values().append(spreadsheetId=sid, range=range_a1, valueInputOption=value_input, insertDataOption="INSERT_ROWS", body={"values": values}).execute()
    updates = result.get("updates", {})
    return {"updatedRange": updates.get("updatedRange"), "updatedRows": updates.get("updatedRows"), "updatedCells": updates.get("updatedCells")}


def _sheets_batch_write(spreadsheet_id, data, value_input="USER_ENTERED"):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    result = _sheets().values().batchUpdate(spreadsheetId=sid, body={"valueInputOption": value_input, "data": data}).execute()
    return {"totalUpdatedRows": result.get("totalUpdatedRows"), "totalUpdatedColumns": result.get("totalUpdatedColumns"), "totalUpdatedCells": result.get("totalUpdatedCells"), "totalUpdatedSheets": result.get("totalUpdatedSheets")}


def _sheets_update_cell(spreadsheet_id, cell_a1, value):
    return _sheets_write_range(spreadsheet_id, cell_a1, [[value]])


def _sheets_clear_range(spreadsheet_id, range_a1):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    result = _sheets().values().clear(spreadsheetId=sid, range=range_a1, body={}).execute()
    return {"clearedRange": result.get("clearedRange")}


def _sheets_insert_rows_columns(spreadsheet_id, sheet_id, dimension, start_index, end_index, inherit_from_before=False):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"insertDimension": {"range": {"sheetId": sheet_id, "dimension": dimension.upper(), "startIndex": start_index, "endIndex": end_index}, "inheritFromBefore": inherit_from_before}}]}).execute()
    return {"inserted": dimension.lower(), "count": end_index - start_index, "startIndex": start_index}


def _sheets_delete_rows_columns(spreadsheet_id, sheet_id, dimension, start_index, end_index):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"deleteDimension": {"range": {"sheetId": sheet_id, "dimension": dimension.upper(), "startIndex": start_index, "endIndex": end_index}}}]}).execute()
    return {"deleted": dimension.lower(), "count": end_index - start_index, "startIndex": start_index}


def _sheets_format_cells(spreadsheet_id, sheet_id, range_a1, **kwargs):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    text_kwargs = {}
    for k in ("bold", "italic", "underline", "strikethrough", "font_size", "font_family", "font_color"):
        if k in kwargs and kwargs[k] is not None:
            text_kwargs[k] = kwargs[k]
    cell_format, fields_str = _make_cell_format(
        bg_color=kwargs.get("bg_color"), h_align=kwargs.get("h_align"), v_align=kwargs.get("v_align"),
        wrap_strategy=kwargs.get("wrap_strategy"), number_format=kwargs.get("number_format"),
        number_format_type=kwargs.get("number_format_type"), **text_kwargs,
    )
    gr = _a1_to_grid_range(range_a1, sheet_id)
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"repeatCell": {"range": gr, "cell": {"userEnteredFormat": cell_format}, "fields": fields_str}}]}).execute()
    return {"formatted": True, "range": range_a1, "fields": fields_str}


def _sheets_merge_cells(spreadsheet_id, sheet_id, range_a1, merge_type="MERGE_ALL"):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    gr = _a1_to_grid_range(range_a1, sheet_id)
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"mergeCells": {"range": gr, "mergeType": merge_type}}]}).execute()
    return {"merged": True, "range": range_a1, "type": merge_type}


def _sheets_unmerge_cells(spreadsheet_id, sheet_id, range_a1):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    gr = _a1_to_grid_range(range_a1, sheet_id)
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"unmergeCells": {"range": gr}}]}).execute()
    return {"unmerged": True, "range": range_a1}


def _sheets_set_column_width(spreadsheet_id, sheet_id, start_col, end_col, width):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"updateDimensionProperties": {"range": {"sheetId": sheet_id, "dimension": "COLUMNS", "startIndex": start_col, "endIndex": end_col}, "properties": {"pixelSize": width}, "fields": "pixelSize"}}]}).execute()
    return {"updated": True, "columns": f"{start_col}-{end_col}", "width": width}


def _sheets_set_row_height(spreadsheet_id, sheet_id, start_row, end_row, height):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"updateDimensionProperties": {"range": {"sheetId": sheet_id, "dimension": "ROWS", "startIndex": start_row, "endIndex": end_row}, "properties": {"pixelSize": height}, "fields": "pixelSize"}}]}).execute()
    return {"updated": True, "rows": f"{start_row}-{end_row}", "height": height}


def _sheets_auto_resize(spreadsheet_id, sheet_id, dimension, start_index=0, end_index=26):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"autoResizeDimensions": {"dimensions": {"sheetId": sheet_id, "dimension": dimension.upper(), "startIndex": start_index, "endIndex": end_index}}}]}).execute()
    return {"autoResized": True, "dimension": dimension, "range": f"{start_index}-{end_index}"}


def _sheets_set_borders(spreadsheet_id, sheet_id, range_a1, style="SOLID", color="#000000", top=True, bottom=True, left=True, right=True, inner_horizontal=False, inner_vertical=False):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    gr = _a1_to_grid_range(range_a1, sheet_id)
    border_style = {"style": style, "colorStyle": {"rgbColor": _make_color(color)}}
    borders = {}
    if top: borders["top"] = border_style
    if bottom: borders["bottom"] = border_style
    if left: borders["left"] = border_style
    if right: borders["right"] = border_style
    if inner_horizontal: borders["innerHorizontal"] = border_style
    if inner_vertical: borders["innerVertical"] = border_style
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"updateBorders": {"range": gr, **borders}}]}).execute()
    return {"borders_set": True, "range": range_a1, "style": style}


def _sheets_add_conditional_format(spreadsheet_id, sheet_id, range_a1, rule_type, values=None, bg_color=None, text_color=None, bold=None, min_color=None, mid_color=None, max_color=None):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    gr = _a1_to_grid_range(range_a1, sheet_id)
    if rule_type == "COLOR_SCALE":
        rule = {"gradientRule": {"minpoint": {"color": _make_color(min_color or "#FFFFFF"), "type": "MIN"}, "maxpoint": {"color": _make_color(max_color or "#FF0000"), "type": "MAX"}}}
        if mid_color:
            rule["gradientRule"]["midpoint"] = {"color": _make_color(mid_color), "type": "PERCENTILE", "value": "50"}
    else:
        fmt = {}
        if bg_color: fmt["backgroundColor"] = _make_color(bg_color)
        if text_color: fmt["textFormat"] = {"foregroundColor": _make_color(text_color)}
        if bold is not None: fmt.setdefault("textFormat", {})["bold"] = bold
        condition = {"type": rule_type}
        if values: condition["values"] = [{"userEnteredValue": v} for v in values]
        rule = {"booleanRule": {"condition": condition, "format": fmt}}
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"addConditionalFormatRule": {"rule": {"ranges": [gr], **rule}, "index": 0}}]}).execute()
    return {"added": True, "ruleType": rule_type, "range": range_a1}


def _sheets_set_data_validation(spreadsheet_id, sheet_id, range_a1, validation_type, values=None, strict=True, input_message=None):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    gr = _a1_to_grid_range(range_a1, sheet_id)
    if validation_type == "BOOLEAN":
        condition = {"type": "BOOLEAN"}
    else:
        condition = {"type": validation_type}
        if values: condition["values"] = [{"userEnteredValue": v} for v in values]
    rule = {"condition": condition, "strict": strict, "showCustomUi": True}
    if input_message: rule["inputMessage"] = input_message
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"setDataValidation": {"range": gr, "rule": rule}}]}).execute()
    return {"validation_set": True, "type": validation_type, "range": range_a1}


def _sheets_create_named_range(spreadsheet_id, sheet_id, name, range_a1):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    gr = _a1_to_grid_range(range_a1, sheet_id)
    result = _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"addNamedRange": {"namedRange": {"name": name, "range": gr}}}]}).execute()
    nr = result["replies"][0]["addNamedRange"]["namedRange"]
    return {"namedRangeId": nr["namedRangeId"], "name": nr["name"]}


def _sheets_add_filter_view(spreadsheet_id, sheet_id, title, range_a1, sort_column=None, sort_order="ASCENDING", filter_column=None, hidden_values=None):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    gr = _a1_to_grid_range(range_a1, sheet_id)
    fv = {"title": title, "range": gr}
    if sort_column is not None:
        fv["sortSpecs"] = [{"dimensionIndex": sort_column, "sortOrder": sort_order.upper()}]
    if filter_column is not None and hidden_values:
        fv["criteria"] = {str(filter_column): {"hiddenValues": hidden_values}}
    result = _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"addFilterView": {"filter": fv}}]}).execute()
    fv_result = result["replies"][0]["addFilterView"]["filter"]
    return {"filterViewId": fv_result["filterViewId"], "title": fv_result["title"]}


def _sheets_set_basic_filter(spreadsheet_id, sheet_id, range_a1):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    gr = _a1_to_grid_range(range_a1, sheet_id)
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"setBasicFilter": {"filter": {"range": gr}}}]}).execute()
    return {"filter_set": True, "range": range_a1}


def _sheets_clear_basic_filter(spreadsheet_id, sheet_id):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"clearBasicFilter": {"sheetId": sheet_id}}]}).execute()
    return {"filter_cleared": True, "sheetId": sheet_id}


def _sheets_create_pivot_table(spreadsheet_id, source_sheet_id, source_range_a1, target_sheet_id, target_row=0, target_col=0, row_fields=None, value_fields=None, column_fields=None, filter_fields=None):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    source_gr = _a1_to_grid_range(source_range_a1, source_sheet_id)
    pivot = {
        "source": source_gr,
        "rows": [{"sourceColumnOffset": f, "showTotals": True, "sortOrder": "ASCENDING"} for f in row_fields],
        "values": [{"sourceColumnOffset": v["sourceColumnOffset"], "summarizeFunction": v.get("summarizeFunction", "SUM")} for v in value_fields],
    }
    if column_fields:
        pivot["columns"] = [{"sourceColumnOffset": f, "showTotals": True} for f in column_fields]
    if filter_fields:
        pivot["filterSpecs"] = [{"filterCriteria": ff.get("filterCriteria", {}), "columnOffsetIndex": ff["sourceColumnOffset"]} for ff in filter_fields]
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"updateCells": {"rows": [{"values": [{"pivotTable": pivot}]}], "start": {"sheetId": target_sheet_id, "rowIndex": target_row, "columnIndex": target_col}, "fields": "pivotTable"}}]}).execute()
    return {"pivotTable_created": True, "targetSheet": target_sheet_id}


def _sheets_protect_range(spreadsheet_id, sheet_id, range_a1=None, description="", warning_only=False, editors=None):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    protected = {"description": description, "warningOnly": warning_only}
    if range_a1:
        protected["range"] = _a1_to_grid_range(range_a1, sheet_id)
    else:
        protected["sheetId"] = sheet_id
    if editors:
        protected["editors"] = {"users": editors}
    result = _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"addProtectedRange": {"protectedRange": protected}}]}).execute()
    pr = result["replies"][0]["addProtectedRange"]["protectedRange"]
    return {"protectedRangeId": pr["protectedRangeId"], "description": pr.get("description"), "warningOnly": pr.get("warningOnly", False)}


def _sheets_set_note(spreadsheet_id, sheet_id, cell_a1, note):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    gr = _a1_to_grid_range(cell_a1, sheet_id)
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"updateCells": {"range": gr, "rows": [{"values": [{"note": note}]}], "fields": "note"}}]}).execute()
    return {"note_set": True, "cell": cell_a1, "note": note}


def _sheets_get_notes(spreadsheet_id, sheet_name):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    sp = _sheets().get(spreadsheetId=sid, ranges=[f"'{sheet_name}'"], fields="sheets.data.rowData.values.note").execute()
    notes = []
    for sheet in sp.get("sheets", []):
        for grid_data in sheet.get("data", []):
            for row_idx, row_data in enumerate(grid_data.get("rowData", [])):
                for col_idx, cell in enumerate(row_data.get("values", [])):
                    note = cell.get("note")
                    if note:
                        notes.append({"cell": f"{_col_index_to_letter(col_idx)}{row_idx + 1}", "note": note})
    return notes


def _sheets_get_formulas(spreadsheet_id, range_a1):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    result = _sheets().values().get(spreadsheetId=sid, range=range_a1, valueRenderOption="FORMULA").execute()
    values = result.get("values", [])
    formulas = []
    for row_idx, row in enumerate(values):
        for col_idx, val in enumerate(row):
            if isinstance(val, str) and val.startswith("="):
                formulas.append({"cell": f"{_col_index_to_letter(col_idx)}{row_idx + 1}", "formula": val})
    return {"range": range_a1, "formulas": formulas, "total": len(formulas)}


def _sheets_sort_range(spreadsheet_id, sheet_id, range_a1, sort_specs):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    gr = _a1_to_grid_range(range_a1, sheet_id)
    _sheets().batchUpdate(spreadsheetId=sid, body={"requests": [{"sortRange": {"range": gr, "sortSpecs": sort_specs}}]}).execute()
    return {"sorted": True, "range": range_a1}


# =============================================================================
# Sheets Drive operations (inlined from drive_service.py)
# =============================================================================

def _sheets_create_spreadsheet(title, sheet_names=None, folder_id=None):
    body = {"properties": {"title": title}}
    if sheet_names:
        body["sheets"] = [{"properties": {"title": name}} for name in sheet_names]
    sp = get_sheets_service().spreadsheets().create(body=body).execute()
    spreadsheet_id = sp["spreadsheetId"]
    if folder_id:
        drive = get_drive_service()
        f = drive.files().get(fileId=spreadsheet_id, fields="parents", supportsAllDrives=True).execute()
        previous_parents = ",".join(f.get("parents", []))
        drive.files().update(fileId=spreadsheet_id, addParents=folder_id, removeParents=previous_parents, supportsAllDrives=True, fields="id,parents").execute()
    sheets_info = [{"sheetId": s["properties"]["sheetId"], "title": s["properties"]["title"]} for s in sp.get("sheets", [])]
    return {"spreadsheetId": spreadsheet_id, "title": title, "url": sp.get("spreadsheetUrl"), "sheets": sheets_info}


def _sheets_copy_spreadsheet(spreadsheet_id, new_title, folder_id=None):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    body = {"name": new_title}
    if folder_id:
        body["parents"] = [folder_id]
    result = get_drive_service().files().copy(fileId=sid, body=body, supportsAllDrives=True).execute()
    return {"spreadsheetId": result["id"], "title": new_title, "url": f"https://docs.google.com/spreadsheets/d/{result['id']}/edit"}


def _sheets_list_spreadsheets(query=None, folder_id=None, max_results=20, include_shared_drives=True):
    q_parts = ["mimeType='application/vnd.google-apps.spreadsheet'", "trashed=false"]
    if query:
        safe_query = query.replace("'", "\\'")
        q_parts.append(f"name contains '{safe_query}'")
    if folder_id:
        q_parts.append(f"'{folder_id}' in parents")
    result = get_drive_service().files().list(
        q=" and ".join(q_parts), pageSize=max_results,
        fields="files(id,name,modifiedTime,owners,shared,webViewLink)",
        orderBy="modifiedTime desc", supportsAllDrives=include_shared_drives,
        includeItemsFromAllDrives=include_shared_drives,
        corpora="allDrives" if include_shared_drives else "user",
    ).execute()
    spreadsheets = []
    for f in result.get("files", []):
        spreadsheets.append({
            "spreadsheetId": f["id"], "name": f["name"],
            "modifiedTime": f.get("modifiedTime"),
            "owner": f.get("owners", [{}])[0].get("emailAddress") if f.get("owners") else None,
            "shared": f.get("shared", False), "url": f.get("webViewLink"),
        })
    return spreadsheets


def _sheets_share_spreadsheet(spreadsheet_id, email, role="reader", send_notification=True, message=None):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    body = {"type": "user", "role": role, "emailAddress": email}
    get_drive_service().permissions().create(fileId=sid, body=body, sendNotificationEmail=send_notification, emailMessage=message, supportsAllDrives=True).execute()
    return {"shared": True, "email": email, "role": role}


def _sheets_get_permissions(spreadsheet_id):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    result = get_drive_service().permissions().list(fileId=sid, fields="permissions(id,emailAddress,role,type,displayName)", supportsAllDrives=True).execute()
    return [{"id": p["id"], "email": p.get("emailAddress"), "name": p.get("displayName"), "role": p["role"], "type": p["type"]} for p in result.get("permissions", [])]


def _sheets_export_spreadsheet(spreadsheet_id, export_format="xlsx", sheet_id=None, download_dir=None):
    sid = _validate_spreadsheet_id(spreadsheet_id)
    mime_map = {"xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "csv": "text/csv", "pdf": "application/pdf", "ods": "application/vnd.oasis.opendocument.spreadsheet", "tsv": "text/tab-separated-values", "html": "text/html"}
    fmt = export_format.lower()
    if fmt not in mime_map:
        raise ValueError(f"Unsupported format: {fmt}. Supported: {', '.join(mime_map.keys())}")
    if fmt in ("csv", "tsv") and sheet_id is not None:
        import httplib2
        url = f"https://docs.google.com/spreadsheets/d/{sid}/export?format={fmt}&gid={sheet_id}"
        creds = _get_creds()
        http = creds.authorize(httplib2.Http())
        response, content = http.request(url)
        if response.status != 200:
            raise Exception(f"Export failed with status {response.status}")
        data = content
    else:
        data = get_drive_service().files().export(fileId=sid, mimeType=mime_map[fmt]).execute()
    info = get_drive_service().files().get(fileId=sid, fields="name", supportsAllDrives=True).execute()
    safe_name = _re.sub(r'[/\\]', '_', info['name'])
    filename = f"{safe_name}.{fmt}"
    default_dir = DOWNLOAD_DIR
    base_dir = Path(download_dir).resolve() if download_dir else default_dir.resolve()
    allowed = Path.home().resolve()
    if not base_dir.is_relative_to(allowed):
        raise ValueError("download_dir must be within home directory")
    base_dir.mkdir(parents=True, exist_ok=True)
    filepath = base_dir / filename
    if isinstance(data, bytes):
        filepath.write_bytes(data)
    else:
        filepath.write_bytes(str(data).encode())
    return {"exported": True, "format": fmt, "path": str(filepath), "size_bytes": filepath.stat().st_size}


def _sheets_duplicate_as_template(spreadsheet_id, new_title, clear_data=True, folder_id=None):
    result = _sheets_copy_spreadsheet(spreadsheet_id, new_title, folder_id)
    new_id = result["spreadsheetId"]
    if clear_data:
        sp = get_sheets_service().spreadsheets().get(spreadsheetId=new_id).execute()
        for sheet in sp.get("sheets", []):
            title = sheet["properties"]["title"]
            try:
                get_sheets_service().spreadsheets().values().clear(spreadsheetId=new_id, range=f"'{title}'", body={}).execute()
            except Exception:
                pass
    return {"spreadsheetId": new_id, "title": new_title, "url": f"https://docs.google.com/spreadsheets/d/{new_id}/edit", "dataCleared": clear_data}


# =============================================================================
# JSON output helper for Sheets tools
# =============================================================================

def _ok(data: Any) -> list[types.TextContent]:
    text = json.dumps(data, indent=2, ensure_ascii=False, default=str)
    return [types.TextContent(type="text", text=text)]


def _err(msg: str) -> list[types.TextContent]:
    return [types.TextContent(type="text", text=f"Error: {msg}")]


# =============================================================================
# Server
# =============================================================================

server = Server("google-workspace")


# =============================================================================
# Tool definitions
# =============================================================================

@server.list_tools()
async def list_tools() -> list[types.Tool]:
    tools = []

    # --- Drive Tools ---
    tools.extend([
        types.Tool(
            name="search_files",
            description="Search files on Google Drive by name, MIME type, date, folder, etc. Returns name, type, size, modified date and link for each result.",
            inputSchema={
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "Drive API query. Examples:\n  name contains 'report'\n  mimeType = 'application/pdf'\n  modifiedTime > '2025-01-01T00:00:00'\n  'FOLDER_ID' in parents\n  sharedWithMe = true\nLeave empty for recent files."},
                    "page_size": {"type": "integer", "description": "Max results (1-100, default 20)", "default": 20},
                    "order_by": {"type": "string", "description": "Sort: modifiedTime desc | name | createdTime desc", "default": "modifiedTime desc"},
                },
                "required": [],
            },
        ),
        types.Tool(
            name="list_folder",
            description="List all files and subfolders inside a Google Drive folder.",
            inputSchema={
                "type": "object",
                "properties": {
                    "folder_id": {"type": "string", "description": "Folder ID. Use 'root' for My Drive root."},
                    "page_size": {"type": "integer", "description": "Max results (1-100, default 50)", "default": 50},
                },
                "required": ["folder_id"],
            },
        ),
        types.Tool(
            name="get_file_info",
            description="Get full metadata of a file: name, type, size, owner, dates, link.",
            inputSchema={"type": "object", "properties": {"file_id": {"type": "string", "description": "Google Drive file ID"}}, "required": ["file_id"]},
        ),
        types.Tool(
            name="read_file_content",
            description="Read and extract the CONTENT of any file from Google Drive.\n\nSupported formats:\n- Google Docs / Sheets / Slides / Drawings\n- PDF (text extraction)\n- Word (.docx) - full text + tables\n- Excel (.xlsx) - all sheets and rows\n- PowerPoint (.pptx) - all slides text\n- Plain text, CSV, JSON, HTML, Markdown\n- ZIP - shows file listing\n- Images (JPEG, PNG, GIF, WebP) - shown visually\n- Video/audio - returns metadata only",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_id": {"type": "string", "description": "Google Drive file ID"},
                    "max_chars": {"type": "integer", "description": "Max chars to return for text (default 10000)", "default": 10000},
                },
                "required": ["file_id"],
            },
        ),
        types.Tool(
            name="download_file",
            description="Download a file from Google Drive to the local downloads/ folder.\nGoogle Docs -> .docx, Sheets -> .xlsx, Slides -> .pptx (auto-converted).",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_id": {"type": "string", "description": "Google Drive file ID"},
                    "save_as": {"type": "string", "description": "Custom filename (optional). Defaults to Drive filename."},
                },
                "required": ["file_id"],
            },
        ),
        types.Tool(
            name="create_folder",
            description="Create a new folder on Google Drive.",
            inputSchema={"type": "object", "properties": {"name": {"type": "string", "description": "Folder name"}, "parent_id": {"type": "string", "description": "Parent folder ID (default: root)"}}, "required": ["name"]},
        ),
        types.Tool(
            name="upload_file",
            description="Upload a local file to Google Drive.",
            inputSchema={"type": "object", "properties": {"local_path": {"type": "string", "description": "Absolute path to local file"}, "name": {"type": "string", "description": "Name on Drive (default: same as local)"}, "parent_id": {"type": "string", "description": "Destination folder ID (default: root)"}}, "required": ["local_path"]},
        ),
        types.Tool(
            name="move_file",
            description="Move a file or folder to a different folder on Google Drive.",
            inputSchema={"type": "object", "properties": {"file_id": {"type": "string", "description": "File or folder ID to move"}, "new_parent_id": {"type": "string", "description": "Destination folder ID"}}, "required": ["file_id", "new_parent_id"]},
        ),
        types.Tool(
            name="rename_file",
            description="Rename a file or folder on Google Drive.",
            inputSchema={"type": "object", "properties": {"file_id": {"type": "string", "description": "File or folder ID"}, "new_name": {"type": "string", "description": "New name"}}, "required": ["file_id", "new_name"]},
        ),
        types.Tool(
            name="copy_file",
            description="Copy a file on Google Drive.",
            inputSchema={"type": "object", "properties": {"file_id": {"type": "string", "description": "Source file ID"}, "new_name": {"type": "string", "description": "Name for the copy (optional)"}, "parent_id": {"type": "string", "description": "Destination folder ID (optional)"}}, "required": ["file_id"]},
        ),
        types.Tool(
            name="delete_file",
            description="Move a file or folder to trash on Google Drive.",
            inputSchema={"type": "object", "properties": {"file_id": {"type": "string", "description": "File or folder ID to trash"}}, "required": ["file_id"]},
        ),
        types.Tool(
            name="share_file",
            description="Share a file or folder with a user by email.",
            inputSchema={"type": "object", "properties": {"file_id": {"type": "string", "description": "File or folder ID"}, "email": {"type": "string", "description": "Email address to share with"}, "role": {"type": "string", "description": "Permission role: reader, commenter, writer, owner", "default": "reader"}}, "required": ["file_id", "email"]},
        ),
        types.Tool(
            name="list_shared_drives",
            description="List all Shared Drives (Team Drives) the user has access to. Use the returned drive ID with list_folder to browse.",
            inputSchema={"type": "object", "properties": {"page_size": {"type": "integer", "description": "Max results (1-100, default 50)", "default": 50}}, "required": []},
        ),
    ])

    # --- Calendar Tools ---
    tools.extend([
        types.Tool(name="cal_list_calendars", description="List all calendars the user has access to (own + subscribed).", inputSchema={"type": "object", "properties": {}}),
        types.Tool(
            name="cal_list_events",
            description="List events from a calendar. Defaults to primary calendar, next 30 days.",
            inputSchema={"type": "object", "properties": {
                "calendar_id": {"type": "string", "description": "Calendar ID (default: 'primary')"},
                "time_min": {"type": "string", "description": "Start of range, ISO 8601 (default: now)"},
                "time_max": {"type": "string", "description": "End of range, ISO 8601 (default: now + 30 days)"},
                "max_results": {"type": "integer", "description": "Max events (default 50, max 2500)"},
                "query": {"type": "string", "description": "Free-text search filter"},
                "show_deleted": {"type": "boolean", "description": "Include cancelled events"},
            }},
        ),
        types.Tool(name="cal_get_event", description="Get full details of a single calendar event.", inputSchema={"type": "object", "required": ["event_id"], "properties": {"event_id": {"type": "string"}, "calendar_id": {"type": "string", "description": "Default: 'primary'"}}}),
        types.Tool(
            name="cal_create_event",
            description="Create a new calendar event. Supports attendees, Zoom link, reminders, recurrence (including specific days like Tue/Thu/Fri), location, description, all-day events.\n\nRECURRENCE EXAMPLES:\n- Every week: RRULE:FREQ=WEEKLY\n- Every Tue, Thu, Fri: RRULE:FREQ=WEEKLY;BYDAY=TU,TH,FR\n- Every weekday: RRULE:FREQ=WEEKLY;BYDAY=MO,TU,WE,TH,FR\n- Every 2 weeks on Mon: RRULE:FREQ=WEEKLY;INTERVAL=2;BYDAY=MO\n- Daily for 10 times: RRULE:FREQ=DAILY;COUNT=10\n- Monthly on 1st: RRULE:FREQ=MONTHLY;BYMONTHDAY=1\n- Until a date: RRULE:FREQ=WEEKLY;BYDAY=MO,WE;UNTIL=20261231T235959Z\nDay codes: MO,TU,WE,TH,FR,SA,SU",
            inputSchema={"type": "object", "required": ["summary", "start", "end"], "properties": {
                "calendar_id": {"type": "string", "description": "Default: 'primary'"},
                "summary": {"type": "string", "description": "Event title"}, "description": {"type": "string"}, "location": {"type": "string"},
                "start": {"type": "string", "description": "ISO 8601 datetime or date for all-day"}, "end": {"type": "string", "description": "ISO 8601 datetime or date for all-day"},
                "time_zone": {"type": "string", "description": "IANA timezone (e.g. Europe/Kyiv)"},
                "attendees": {"type": "string", "description": "Comma-separated emails of attendees"},
                "recurrence": {"type": "string", "description": "RRULE string"},
                "zoom_link": {"type": "string", "description": "Zoom meeting URL"},
                "reminders_minutes": {"type": "string", "description": "Comma-separated minutes before event for popup reminders"},
                "color_id": {"type": "string", "description": "Event color ID (1-11)"},
                "visibility": {"type": "string", "description": "'default', 'public', 'private', 'confidential'"},
                "send_updates": {"type": "string", "description": "'all', 'externalOnly', 'none' (default: 'none')"},
            }},
        ),
        types.Tool(
            name="cal_update_event",
            description="Update an existing calendar event (any field: title, time, attendees, Zoom link, description, color, recurrence, etc.).",
            inputSchema={"type": "object", "required": ["event_id"], "properties": {
                "event_id": {"type": "string"}, "calendar_id": {"type": "string", "description": "Default: 'primary'"},
                "summary": {"type": "string"}, "description": {"type": "string"}, "location": {"type": "string"},
                "start": {"type": "string"}, "end": {"type": "string"}, "time_zone": {"type": "string"},
                "attendees": {"type": "string", "description": "Comma-separated emails (replaces existing)"},
                "zoom_link": {"type": "string"}, "recurrence": {"type": "string"},
                "color_id": {"type": "string"}, "status": {"type": "string", "description": "'confirmed', 'tentative', 'cancelled'"},
                "send_updates": {"type": "string", "description": "'all', 'externalOnly', 'none'"},
            }},
        ),
        types.Tool(name="cal_delete_event", description="Delete a calendar event.", inputSchema={"type": "object", "required": ["event_id"], "properties": {"event_id": {"type": "string"}, "calendar_id": {"type": "string", "description": "Default: 'primary'"}, "send_updates": {"type": "string", "description": "'all', 'externalOnly', 'none'"}}}),
        types.Tool(name="cal_move_event", description="Move an event from one calendar to another.", inputSchema={"type": "object", "required": ["event_id", "destination_calendar_id"], "properties": {"event_id": {"type": "string"}, "calendar_id": {"type": "string", "description": "Source calendar (default: 'primary')"}, "destination_calendar_id": {"type": "string"}}}),
        types.Tool(name="cal_quick_add", description="Create an event from a natural language string (e.g. 'Meeting with Bob tomorrow at 3pm for 1 hour').", inputSchema={"type": "object", "required": ["text"], "properties": {"text": {"type": "string", "description": "Natural language event description"}, "calendar_id": {"type": "string", "description": "Default: 'primary'"}}}),
        types.Tool(name="cal_freebusy", description="Check availability (free/busy) for one or more calendars in a time range.", inputSchema={"type": "object", "required": ["time_min", "time_max"], "properties": {"time_min": {"type": "string", "description": "Start of range, ISO 8601"}, "time_max": {"type": "string", "description": "End of range, ISO 8601"}, "calendar_ids": {"type": "string", "description": "Comma-separated calendar IDs (default: 'primary')"}}}),
        types.Tool(name="cal_get_colors", description="Get all available calendar and event color definitions.", inputSchema={"type": "object", "properties": {}}),
        types.Tool(name="cal_list_recurring_instances", description="List all instances of a recurring event.", inputSchema={"type": "object", "required": ["event_id"], "properties": {"event_id": {"type": "string"}, "calendar_id": {"type": "string", "description": "Default: 'primary'"}, "time_min": {"type": "string"}, "time_max": {"type": "string"}, "max_results": {"type": "integer"}}}),
        types.Tool(name="cal_respond_to_event", description="Set your RSVP response to an event (accepted, declined, tentative).", inputSchema={"type": "object", "required": ["event_id", "response"], "properties": {"event_id": {"type": "string"}, "calendar_id": {"type": "string", "description": "Default: 'primary'"}, "response": {"type": "string", "description": "'accepted', 'declined', 'tentative'"}}}),
    ])

    # --- Sheets Tools ---
    tools.extend([
        # File operations
        types.Tool(name="create_spreadsheet", description="Create a new Google Sheets spreadsheet. Optionally specify sheet names and target folder.", inputSchema={"type": "object", "required": ["title"], "properties": {"title": {"type": "string", "description": "Spreadsheet title"}, "sheet_names": {"type": "array", "items": {"type": "string"}, "description": "List of sheet names to create"}, "folder_id": {"type": "string", "description": "Google Drive folder ID to place the spreadsheet in"}}}),
        types.Tool(name="copy_spreadsheet", description="Copy an existing spreadsheet to a new one.", inputSchema={"type": "object", "required": ["spreadsheet_id", "new_title"], "properties": {"spreadsheet_id": {"type": "string", "description": "Spreadsheet ID or URL"}, "new_title": {"type": "string", "description": "Title for the copy"}, "folder_id": {"type": "string", "description": "Folder ID for the copy"}}}),
        types.Tool(name="get_spreadsheet_info", description="Get full metadata: title, locale, timezone, all sheets with sizes, named ranges.", inputSchema={"type": "object", "required": ["spreadsheet_id"], "properties": {"spreadsheet_id": {"type": "string", "description": "Spreadsheet ID or URL"}}}),
        types.Tool(name="list_spreadsheets", description="Search for spreadsheets on Google Drive by name, folder, including Shared Drives.", inputSchema={"type": "object", "properties": {"query": {"type": "string", "description": "Search query (name contains)"}, "folder_id": {"type": "string", "description": "Limit search to specific folder"}, "max_results": {"type": "integer", "description": "Max results (default 20)", "default": 20}, "include_shared_drives": {"type": "boolean", "description": "Include Shared Drives (default true)", "default": True}}}),
        types.Tool(name="share_spreadsheet", description="Share a spreadsheet with a user by email. Roles: reader, commenter, writer, owner.", inputSchema={"type": "object", "required": ["spreadsheet_id", "email", "role"], "properties": {"spreadsheet_id": {"type": "string", "description": "Spreadsheet ID or URL"}, "email": {"type": "string", "description": "Email address"}, "role": {"type": "string", "enum": ["reader", "commenter", "writer", "owner"]}, "send_notification": {"type": "boolean", "default": True}, "message": {"type": "string"}}}),
        types.Tool(name="get_permissions", description="List all users who have access to a spreadsheet with their roles.", inputSchema={"type": "object", "required": ["spreadsheet_id"], "properties": {"spreadsheet_id": {"type": "string", "description": "Spreadsheet ID or URL"}}}),
        # Sheet management
        types.Tool(name="list_sheets", description="List all sheets (tabs) in a spreadsheet with their IDs, sizes, and properties.", inputSchema={"type": "object", "required": ["spreadsheet_id"], "properties": {"spreadsheet_id": {"type": "string", "description": "Spreadsheet ID or URL"}}}),
        types.Tool(name="add_sheet", description="Add a new sheet (tab) to a spreadsheet.", inputSchema={"type": "object", "required": ["spreadsheet_id", "title"], "properties": {"spreadsheet_id": {"type": "string"}, "title": {"type": "string"}, "rows": {"type": "integer", "default": 1000}, "columns": {"type": "integer", "default": 26}, "tab_color": {"type": "string", "description": "Tab color in hex (#RRGGBB)"}}}),
        types.Tool(name="delete_sheet", description="Delete a sheet (tab) from a spreadsheet.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer", "description": "Sheet ID (numeric)"}}}),
        types.Tool(name="rename_sheet", description="Rename a sheet (tab) in a spreadsheet.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "new_title"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "new_title": {"type": "string"}}}),
        types.Tool(name="copy_sheet", description="Copy a sheet to the same or another spreadsheet.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "destination_spreadsheet_id": {"type": "string", "description": "Target spreadsheet ID (default: same)"}}}),
        types.Tool(name="reorder_sheets", description="Change the order of sheets (tabs) in a spreadsheet.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_order"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_order": {"type": "array", "items": {"type": "object", "required": ["sheetId", "index"], "properties": {"sheetId": {"type": "integer"}, "index": {"type": "integer"}}}}}}),
        # Reading data
        types.Tool(name="read_range", description="Read data from a range (A1 notation). Supports named ranges. Returns values as 2D array.", inputSchema={"type": "object", "required": ["spreadsheet_id", "range"], "properties": {"spreadsheet_id": {"type": "string"}, "range": {"type": "string", "description": "Range in A1 notation (e.g. 'Sheet1!A1:D10')"}, "value_render": {"type": "string", "enum": ["FORMATTED_VALUE", "UNFORMATTED_VALUE", "FORMULA"]}}}),
        types.Tool(name="batch_read", description="Read multiple ranges in a single API call.", inputSchema={"type": "object", "required": ["spreadsheet_id", "ranges"], "properties": {"spreadsheet_id": {"type": "string"}, "ranges": {"type": "array", "items": {"type": "string"}}, "value_render": {"type": "string", "enum": ["FORMATTED_VALUE", "UNFORMATTED_VALUE", "FORMULA"]}}}),
        types.Tool(name="get_cell", description="Get a single cell's value, formula, and format info.", inputSchema={"type": "object", "required": ["spreadsheet_id", "cell"], "properties": {"spreadsheet_id": {"type": "string"}, "cell": {"type": "string", "description": "Cell in A1 notation (e.g. 'Sheet1!B5')"}}}),
        types.Tool(name="search_values", description="Search for a value across all sheets or a specific sheet. Case-insensitive partial match.", inputSchema={"type": "object", "required": ["spreadsheet_id", "query"], "properties": {"spreadsheet_id": {"type": "string"}, "query": {"type": "string", "description": "Text to search for"}, "sheet_name": {"type": "string"}}}),
        types.Tool(name="get_sheet_as_table", description="Get entire sheet as a structured table with headers and row objects.", inputSchema={"type": "object", "required": ["spreadsheet_id"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_name": {"type": "string"}, "has_header": {"type": "boolean", "default": True}, "max_rows": {"type": "integer", "default": 0}}}),
        # Writing data
        types.Tool(name="write_range", description="Write data to a range (overwrites existing data). Values as 2D array.", inputSchema={"type": "object", "required": ["spreadsheet_id", "range", "values"], "properties": {"spreadsheet_id": {"type": "string"}, "range": {"type": "string"}, "values": {"type": "array", "items": {"type": "array"}, "description": "2D array of values"}, "value_input": {"type": "string", "enum": ["USER_ENTERED", "RAW"]}}}),
        types.Tool(name="append_rows", description="Append rows after the last row with data.", inputSchema={"type": "object", "required": ["spreadsheet_id", "range", "values"], "properties": {"spreadsheet_id": {"type": "string"}, "range": {"type": "string"}, "values": {"type": "array", "items": {"type": "array"}}, "value_input": {"type": "string", "enum": ["USER_ENTERED", "RAW"], "default": "USER_ENTERED"}}}),
        types.Tool(name="batch_write", description="Write to multiple ranges in a single API call.", inputSchema={"type": "object", "required": ["spreadsheet_id", "data"], "properties": {"spreadsheet_id": {"type": "string"}, "data": {"type": "array", "items": {"type": "object", "required": ["range", "values"], "properties": {"range": {"type": "string"}, "values": {"type": "array", "items": {"type": "array"}}}}}, "value_input": {"type": "string", "enum": ["USER_ENTERED", "RAW"], "default": "USER_ENTERED"}}}),
        types.Tool(name="update_cell", description="Update a single cell value or formula.", inputSchema={"type": "object", "required": ["spreadsheet_id", "cell", "value"], "properties": {"spreadsheet_id": {"type": "string"}, "cell": {"type": "string"}, "value": {"description": "Value to write"}}}),
        types.Tool(name="clear_range", description="Clear all values in a range (keeps formatting).", inputSchema={"type": "object", "required": ["spreadsheet_id", "range"], "properties": {"spreadsheet_id": {"type": "string"}, "range": {"type": "string"}}}),
        types.Tool(name="insert_rows_columns", description="Insert empty rows or columns at a specific position.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "dimension", "start_index", "end_index"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "dimension": {"type": "string", "enum": ["ROWS", "COLUMNS"]}, "start_index": {"type": "integer"}, "end_index": {"type": "integer"}, "inherit_from_before": {"type": "boolean", "default": False}}}),
        types.Tool(name="delete_rows_columns", description="Delete rows or columns from a sheet.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "dimension", "start_index", "end_index"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "dimension": {"type": "string", "enum": ["ROWS", "COLUMNS"]}, "start_index": {"type": "integer"}, "end_index": {"type": "integer"}}}),
        # Formatting
        types.Tool(name="format_cells", description="Format cells: font, size, color, bold, italic, alignment, borders, number format, background color.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "range"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "range": {"type": "string"}, "bg_color": {"type": "string"}, "font_color": {"type": "string"}, "bold": {"type": "boolean"}, "italic": {"type": "boolean"}, "underline": {"type": "boolean"}, "strikethrough": {"type": "boolean"}, "font_size": {"type": "integer"}, "font_family": {"type": "string"}, "h_align": {"type": "string", "enum": ["LEFT", "CENTER", "RIGHT"]}, "v_align": {"type": "string", "enum": ["TOP", "MIDDLE", "BOTTOM"]}, "wrap_strategy": {"type": "string", "enum": ["OVERFLOW_CELL", "CLIP", "WRAP"]}, "number_format": {"type": "string"}, "number_format_type": {"type": "string", "enum": ["NUMBER", "CURRENCY", "PERCENT", "DATE", "TIME", "DATE_TIME", "SCIENTIFIC", "TEXT"]}}}),
        types.Tool(name="merge_cells", description="Merge cells in a range.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "range"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "range": {"type": "string"}, "merge_type": {"type": "string", "enum": ["MERGE_ALL", "MERGE_COLUMNS", "MERGE_ROWS"]}}}),
        types.Tool(name="unmerge_cells", description="Unmerge previously merged cells.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "range"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "range": {"type": "string"}}}),
        types.Tool(name="set_column_width", description="Set width of one or more columns in pixels.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "start_col", "end_col", "width"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "start_col": {"type": "integer"}, "end_col": {"type": "integer"}, "width": {"type": "integer"}}}),
        types.Tool(name="set_row_height", description="Set height of one or more rows in pixels.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "start_row", "end_row", "height"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "start_row": {"type": "integer"}, "end_row": {"type": "integer"}, "height": {"type": "integer"}}}),
        types.Tool(name="auto_resize", description="Auto-resize columns or rows to fit content.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "dimension"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "dimension": {"type": "string", "enum": ["COLUMNS", "ROWS"]}, "start_index": {"type": "integer", "default": 0}, "end_index": {"type": "integer", "default": 26}}}),
        types.Tool(name="set_borders", description="Set borders on a range of cells.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "range"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "range": {"type": "string"}, "style": {"type": "string", "enum": ["SOLID", "SOLID_MEDIUM", "SOLID_THICK", "DASHED", "DOTTED", "DOUBLE", "NONE"]}, "color": {"type": "string"}, "top": {"type": "boolean", "default": True}, "bottom": {"type": "boolean", "default": True}, "left": {"type": "boolean", "default": True}, "right": {"type": "boolean", "default": True}, "inner_horizontal": {"type": "boolean", "default": False}, "inner_vertical": {"type": "boolean", "default": False}}}),
        # Advanced features
        types.Tool(name="add_conditional_format", description="Add conditional formatting rule.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "range", "rule_type"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "range": {"type": "string"}, "rule_type": {"type": "string", "enum": ["NUMBER_GREATER", "NUMBER_LESS", "NUMBER_BETWEEN", "TEXT_CONTAINS", "TEXT_NOT_CONTAINS", "CUSTOM_FORMULA", "NOT_BLANK", "BLANK", "COLOR_SCALE"]}, "values": {"type": "array", "items": {"type": "string"}}, "bg_color": {"type": "string"}, "text_color": {"type": "string"}, "bold": {"type": "boolean"}, "min_color": {"type": "string"}, "mid_color": {"type": "string"}, "max_color": {"type": "string"}}}),
        types.Tool(name="set_data_validation", description="Set data validation on cells: dropdown lists, number ranges, dates, checkboxes, custom formulas.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "range", "validation_type"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "range": {"type": "string"}, "validation_type": {"type": "string", "enum": ["ONE_OF_LIST", "NUMBER_BETWEEN", "NUMBER_GREATER", "NUMBER_LESS", "DATE_AFTER", "DATE_BEFORE", "CUSTOM_FORMULA", "BOOLEAN"]}, "values": {"type": "array", "items": {"type": "string"}}, "strict": {"type": "boolean", "default": True}, "input_message": {"type": "string"}}}),
        types.Tool(name="create_named_range", description="Create a named range for easier reference in formulas.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "name", "range"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "name": {"type": "string"}, "range": {"type": "string"}}}),
        types.Tool(name="add_filter_view", description="Create a saved filter view.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "title", "range"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "title": {"type": "string"}, "range": {"type": "string"}, "sort_column": {"type": "integer"}, "sort_order": {"type": "string", "enum": ["ASCENDING", "DESCENDING"]}, "filter_column": {"type": "integer"}, "hidden_values": {"type": "array", "items": {"type": "string"}}}}),
        types.Tool(name="set_basic_filter", description="Set or clear a basic auto-filter on a sheet.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "range": {"type": "string", "description": "Range for filter. If omitted, clears existing filter."}}}),
        types.Tool(name="create_pivot_table", description="Create a pivot table from source data.", inputSchema={"type": "object", "required": ["spreadsheet_id", "source_sheet_id", "source_range", "target_sheet_id", "row_fields", "value_fields"], "properties": {"spreadsheet_id": {"type": "string"}, "source_sheet_id": {"type": "integer"}, "source_range": {"type": "string"}, "target_sheet_id": {"type": "integer"}, "target_row": {"type": "integer", "default": 0}, "target_col": {"type": "integer", "default": 0}, "row_fields": {"type": "array", "items": {"type": "integer"}}, "value_fields": {"type": "array", "items": {"type": "object", "required": ["sourceColumnOffset"], "properties": {"sourceColumnOffset": {"type": "integer"}, "summarizeFunction": {"type": "string", "enum": ["SUM", "COUNT", "AVERAGE", "MAX", "MIN", "COUNTA", "COUNTUNIQUE", "MEDIAN"], "default": "SUM"}}}}, "column_fields": {"type": "array", "items": {"type": "integer"}}}}),
        # Protection & notes
        types.Tool(name="protect_range", description="Protect a range or entire sheet from editing.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "range": {"type": "string"}, "description": {"type": "string"}, "warning_only": {"type": "boolean", "default": False}, "editors": {"type": "array", "items": {"type": "string"}}}}),
        types.Tool(name="set_note", description="Add, update, or clear a note on a cell.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "cell", "note"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "cell": {"type": "string"}, "note": {"type": "string"}}}),
        types.Tool(name="get_notes", description="Get all notes from a specific sheet.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_name"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_name": {"type": "string"}}}),
        # Export & utilities
        types.Tool(name="export_spreadsheet", description="Export/download spreadsheet as XLSX, CSV, PDF, ODS, TSV, or HTML.", inputSchema={"type": "object", "required": ["spreadsheet_id", "format"], "properties": {"spreadsheet_id": {"type": "string"}, "format": {"type": "string", "enum": ["xlsx", "csv", "pdf", "ods", "tsv", "html"]}, "sheet_id": {"type": "integer"}, "download_dir": {"type": "string"}}}),
        types.Tool(name="get_formulas", description="Show all formulas in a range.", inputSchema={"type": "object", "required": ["spreadsheet_id", "range"], "properties": {"spreadsheet_id": {"type": "string"}, "range": {"type": "string"}}}),
        types.Tool(name="duplicate_as_template", description="Copy a spreadsheet as a template -- keeps formatting but clears all data.", inputSchema={"type": "object", "required": ["spreadsheet_id", "new_title"], "properties": {"spreadsheet_id": {"type": "string"}, "new_title": {"type": "string"}, "folder_id": {"type": "string"}, "clear_data": {"type": "boolean", "default": True}}}),
        types.Tool(name="sort_range", description="Sort a range of data by one or more columns.", inputSchema={"type": "object", "required": ["spreadsheet_id", "sheet_id", "range", "sort_specs"], "properties": {"spreadsheet_id": {"type": "string"}, "sheet_id": {"type": "integer"}, "range": {"type": "string"}, "sort_specs": {"type": "array", "items": {"type": "object", "required": ["dimensionIndex"], "properties": {"dimensionIndex": {"type": "integer"}, "sortOrder": {"type": "string", "enum": ["ASCENDING", "DESCENDING"], "default": "ASCENDING"}}}}}}),
    ])

    # --- TODO: Gmail Tools (requires gmail.modify scope) ---
    # types.Tool(name="gmail_send_email", ...),
    # types.Tool(name="gmail_list_emails", ...),
    # types.Tool(name="gmail_read_email", ...),

    # --- TODO: Google Docs Tools (requires documents scope) ---
    # types.Tool(name="docs_create_doc", ...),
    # types.Tool(name="docs_read_doc", ...),
    # types.Tool(name="docs_update_doc", ...),

    # --- TODO: Google Tasks Tools (requires tasks scope) ---
    # types.Tool(name="tasks_list_tasks", ...),
    # types.Tool(name="tasks_create_task", ...),
    # types.Tool(name="tasks_complete_task", ...),

    return tools


# =============================================================================
# Tool dispatch
# =============================================================================

@server.call_tool()
async def call_tool(name: str, arguments: dict[str, Any]):
    try:
        return await asyncio.to_thread(_call_tool_inner, name, arguments)
    except HttpError as e:
        status = e.resp.status if hasattr(e, 'resp') else 'unknown'
        return [types.TextContent(type="text", text=f"Error: Google API returned HTTP {status}: {_safe_error(e)}")]
    except Exception as e:
        return [types.TextContent(type="text", text=f"Error: {_format_error(e)}")]


def _call_tool_inner(name: str, arguments: dict[str, Any]):
    # --- Calendar dispatch ---
    if name.startswith("cal_"):
        return _cal_dispatch(name, arguments)

    # --- Sheets dispatch ---
    if name in _SHEETS_TOOLS:
        return _sheets_dispatch(name, arguments)

    # --- Drive dispatch ---
    return _drive_dispatch(name, arguments)


# =============================================================================
# Drive dispatch
# =============================================================================

def _drive_dispatch(name: str, arguments: dict[str, Any]):
    service = get_drive_service()

    if name == "search_files":
        q = arguments.get("query", "").strip()
        if not q:
            q = "trashed = false"
        elif "trashed" not in q:
            q = f"({q}) and trashed = false"
        page_size = min(max(int(arguments.get("page_size", 20)), 1), 100)
        order_by = arguments.get("order_by", "modifiedTime desc")
        results = service.files().list(
            q=q, pageSize=page_size, orderBy=order_by,
            fields="files(id,name,mimeType,size,modifiedTime,driveId,teamDriveId)",
            supportsAllDrives=True, includeItemsFromAllDrives=True,
        ).execute()
        files = results.get("files", [])
        if not files:
            return [types.TextContent(type="text", text="No files found.")]
        rows = [fmt_file(f) for f in files]
        out = f"Found **{len(rows)}** file(s):\n\n"
        out += "| # | Name | Type | Size | Modified | Link |\n|---|------|------|------|----------|------|\n"
        for i, r in enumerate(rows, 1):
            out += f"| {i} | {r['name']} | {r['type']} | {r['size']} | {r['modified']} | [Open]({r['link']}) |\n"
        out += "\n**File IDs** (for read_file_content / download_file):\n"
        out += "\n".join(f"  `{r['id']}` — {r['name']}" for r in rows)
        return [types.TextContent(type="text", text=out)]

    if name == "list_folder":
        folder_id = arguments["folder_id"]
        page_size = min(max(int(arguments.get("page_size", 50)), 1), 100)
        results = service.files().list(
            q=f"'{folder_id}' in parents and trashed = false",
            pageSize=page_size, orderBy="folder,name",
            fields="files(id,name,mimeType,size,modifiedTime,driveId)",
            supportsAllDrives=True, includeItemsFromAllDrives=True,
        ).execute()
        files = results.get("files", [])
        if not files:
            return [types.TextContent(type="text", text="Folder is empty or ID not found.")]
        rows = [fmt_file(f) for f in files]
        out = f"Folder `{folder_id}` — **{len(rows)}** item(s):\n\n"
        out += "| # | Name | Type | Size | Modified |\n|---|------|------|------|----------|\n"
        for i, r in enumerate(rows, 1):
            out += f"| {i} | [{r['name']}]({r['link']}) | {r['type']} | {r['size']} | {r['modified']} |\n"
        out += "\n**File IDs:**\n"
        out += "\n".join(f"  `{r['id']}` — {r['name']}" for r in rows)
        return [types.TextContent(type="text", text=out)]

    if name == "list_shared_drives":
        page_size = min(max(int(arguments.get("page_size", 50)), 1), 100)
        results = service.drives().list(pageSize=page_size, fields="drives(id,name)").execute()
        drives = results.get("drives", [])
        if not drives:
            return [types.TextContent(type="text", text="No shared drives found.")]
        out = f"Found **{len(drives)}** shared drive(s):\n\n| # | Name | ID (use as folder_id) |\n|---|------|-----------------------|\n"
        for i, d in enumerate(drives, 1):
            out += f"| {i} | {d['name']} | `{d['id']}` |\n"
        out += "\nUse the drive ID with `list_folder` to browse contents."
        return [types.TextContent(type="text", text=out)]

    if name == "get_file_info":
        file_id = arguments["file_id"]
        f = service.files().get(
            fileId=file_id,
            fields="id,name,mimeType,size,createdTime,modifiedTime,owners,shared,webViewLink,parents,description,imageMediaMetadata,videoMediaMetadata,driveId",
            supportsAllDrives=True,
        ).execute()
        mime = f.get("mimeType", "")
        owners = ", ".join(o.get("displayName", o.get("emailAddress", "?")) for o in f.get("owners", []))
        out = (
            f"**{f.get('name')}**\n\n"
            f"- **ID:** `{f['id']}`\n- **Type:** {MIME_LABELS.get(mime, mime)}\n"
            f"- **Size:** {fmt_size(f.get('size'))}\n- **Created:** {f.get('createdTime', '')[:10]}\n"
            f"- **Modified:** {f.get('modifiedTime', '')[:10]}\n- **Owner(s):** {owners}\n"
            f"- **Shared:** {f.get('shared', False)}\n- **Link:** {f.get('webViewLink') or file_link(file_id, mime)}\n"
        )
        if f.get("description"): out += f"- **Description:** {f['description']}\n"
        if f.get("imageMediaMetadata"):
            m = f["imageMediaMetadata"]
            out += f"- **Image dimensions:** {m.get('width')}x{m.get('height')} px\n"
        if f.get("videoMediaMetadata"):
            m = f["videoMediaMetadata"]
            dur = m.get("durationMillis")
            dur_s = f"{int(dur)//1000}s" if dur else "?"
            out += f"- **Video:** {m.get('width')}x{m.get('height')} px, {dur_s}\n"
        return [types.TextContent(type="text", text=out)]

    if name == "read_file_content":
        file_id = arguments["file_id"]
        max_chars = int(arguments.get("max_chars", 10000))
        meta = service.files().get(fileId=file_id, fields="name,mimeType,size,imageMediaMetadata,videoMediaMetadata", supportsAllDrives=True).execute()
        mime = meta.get("mimeType", "")
        fname = meta.get("name", file_id)

        if mime == "application/vnd.google-apps.folder":
            return [types.TextContent(type="text", text=f"'{fname}' is a folder. Use list_folder with its ID.")]

        if mime.startswith("image/") and mime != "image/svg+xml":
            raw = download_bytes(service, file_id, mime)
            try:
                from PIL import Image
                img = Image.open(io.BytesIO(raw))
                w, h = img.size
                if w > 1200 or h > 1200:
                    img.thumbnail((1200, 1200), Image.LANCZOS)
                out_buf = io.BytesIO()
                fmt = "JPEG" if mime == "image/jpeg" else "PNG"
                img.save(out_buf, format=fmt, quality=85)
                raw = out_buf.getvalue()
                mime_out = "image/jpeg" if fmt == "JPEG" else "image/png"
            except Exception:
                mime_out = mime
            b64 = base64.standard_b64encode(raw).decode()
            m = meta.get("imageMediaMetadata", {})
            caption = f"**{fname}**"
            if m: caption += f"  ({m.get('width')}x{m.get('height')} px)"
            return [types.TextContent(type="text", text=caption), types.ImageContent(type="image", data=b64, mimeType=mime_out)]

        if mime == "image/svg+xml":
            raw = download_bytes(service, file_id, mime)
            return [types.TextContent(type="text", text=f"**{fname}** (SVG)\n\n```xml\n{truncate(raw.decode('utf-8', errors='replace'), max_chars)}\n```")]

        if mime.startswith("video/") or mime.startswith("audio/"):
            vm = meta.get("videoMediaMetadata", {})
            dur = vm.get("durationMillis")
            dur_str = f"{int(dur)//1000//60}m {int(dur)//1000%60}s" if dur else "unknown"
            info = f"**{fname}** ({MIME_LABELS.get(mime, mime)})\n\nCannot display video/audio directly.\n- Size: {fmt_size(meta.get('size'))}\n"
            if vm: info += f"- Resolution: {vm.get('width')}x{vm.get('height')}\n- Duration: {dur_str}\n"
            info += f"\nUse `download_file` to save it locally, then open with a media player."
            return [types.TextContent(type="text", text=info)]

        if mime.startswith("application/vnd.google-apps."):
            raw = export_bytes(service, file_id, "text/plain")
            text = raw.decode("utf-8", errors="replace")
            label = MIME_LABELS.get(mime, mime)
            return [types.TextContent(type="text", text=f"**{fname}** ({label})\n\n{truncate(text, max_chars)}")]

        if mime == "application/pdf":
            raw = download_bytes(service, file_id, mime)
            text = read_pdf(raw)
            return [types.TextContent(type="text", text=f"**{fname}** (PDF)\n\n{truncate(text, max_chars)}")]

        if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            raw = download_bytes(service, file_id, mime)
            text = read_docx(raw)
            return [types.TextContent(type="text", text=f"**{fname}** (Word)\n\n{truncate(text, max_chars)}")]

        if mime in ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel"):
            raw = download_bytes(service, file_id, mime)
            text = read_xlsx(raw)
            return [types.TextContent(type="text", text=f"**{fname}** (Excel)\n\n{truncate(text, max_chars)}")]

        if mime in ("application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"):
            raw = download_bytes(service, file_id, mime)
            text = read_pptx(raw)
            return [types.TextContent(type="text", text=f"**{fname}** (PowerPoint)\n\n{truncate(text, max_chars)}")]

        if mime in ("application/zip", "application/x-zip-compressed"):
            raw = download_bytes(service, file_id, mime)
            text = read_zip(raw)
            return [types.TextContent(type="text", text=f"**{fname}**\n\n{text}")]

        if mime.startswith("text/") or mime in ("application/json",):
            raw = download_bytes(service, file_id, mime)
            text = raw.decode("utf-8", errors="replace")
            ext = Path(fname).suffix.lower()
            lang = {".json": "json", ".html": "html", ".htm": "html", ".md": "markdown", ".csv": "csv", ".py": "python", ".js": "javascript", ".ts": "typescript"}.get(ext, "")
            return [types.TextContent(type="text", text=f"**{fname}**\n\n```{lang}\n{truncate(text, max_chars)}\n```")]

        size = int(meta.get("size") or 0)
        if size < 5 * 1024 * 1024:
            try:
                raw = download_bytes(service, file_id, mime)
                text = raw.decode("utf-8", errors="replace")
                return [types.TextContent(type="text", text=f"**{fname}** ({MIME_LABELS.get(mime, mime)})\n\n```\n{truncate(text, max_chars)}\n```")]
            except Exception:
                pass
        return [types.TextContent(type="text", text=f"**{fname}** ({MIME_LABELS.get(mime, mime)})\n\nCannot display this file type inline. Use `download_file` to save it locally.")]

    if name == "download_file":
        file_id = arguments["file_id"]
        save_as = arguments.get("save_as", "").strip()
        meta = service.files().get(fileId=file_id, fields="name,mimeType,size", supportsAllDrives=True).execute()
        mime = meta.get("mimeType", "")
        orig_name = meta.get("name", file_id)
        file_size = int(meta.get("size") or 0)
        if file_size > MAX_DOWNLOAD_SIZE:
            return [types.TextContent(type="text", text=f"Error: file too large ({fmt_size(file_size)}). Max allowed: {fmt_size(MAX_DOWNLOAD_SIZE)}")]
        DOWNLOAD_DIR.mkdir(exist_ok=True)
        if mime in EXPORT_FORMATS:
            export_mime, ext = EXPORT_FORMATS[mime]
            filename = save_as or (orig_name + ext)
            request = service.files().export_media(fileId=file_id, mimeType=export_mime)
        elif mime.startswith("application/vnd.google-apps."):
            return [types.TextContent(type="text", text=f"'{orig_name}' ({mime.split('.')[-1]}) — no export format available for this type.")]
        else:
            filename = save_as or orig_name
            request = service.files().get_media(fileId=file_id)
        dest = (DOWNLOAD_DIR / filename).resolve()
        if not dest.is_relative_to(DOWNLOAD_DIR.resolve()):
            return [types.TextContent(type="text", text="Error: invalid filename — path traversal detected")]
        try:
            with open(dest, "wb") as fh:
                dl = MediaIoBaseDownload(fh, request)
                done = False
                while not done:
                    _, done = dl.next_chunk()
        except Exception:
            if dest.exists(): dest.unlink()
            raise
        size = dest.stat().st_size
        return [types.TextContent(type="text", text=f"Saved **{filename}** ({fmt_size(size)}) -> `{dest}`")]

    if name == "create_folder":
        folder_name = arguments["name"]
        parent_id = arguments.get("parent_id", "root")
        metadata = {"name": folder_name, "mimeType": "application/vnd.google-apps.folder", "parents": [parent_id]}
        folder = service.files().create(body=metadata, fields="id,name,webViewLink", supportsAllDrives=True).execute()
        return [types.TextContent(type="text", text=f"Created folder **{folder['name']}**\n- ID: `{folder['id']}`\n- Link: {folder.get('webViewLink', '')}")]

    if name == "upload_file":
        from googleapiclient.http import MediaFileUpload
        local_path = Path(arguments["local_path"]).resolve()
        home = Path.home().resolve()
        if not local_path.is_relative_to(home):
            return [types.TextContent(type="text", text="Error: upload restricted to files within home directory")]
        if not local_path.exists():
            return [types.TextContent(type="text", text=f"File not found: {local_path}")]
        file_name = arguments.get("name", local_path.name)
        parent_id = arguments.get("parent_id", "root")
        metadata = {"name": file_name, "parents": [parent_id]}
        import mimetypes
        mime_type = mimetypes.guess_type(str(local_path))[0] or "application/octet-stream"
        media = MediaFileUpload(str(local_path), mimetype=mime_type, resumable=True)
        f = service.files().create(body=metadata, media_body=media, fields="id,name,size,webViewLink", supportsAllDrives=True).execute()
        return [types.TextContent(type="text", text=f"Uploaded **{f['name']}** ({fmt_size(f.get('size'))})\n- ID: `{f['id']}`\n- Link: {f.get('webViewLink', '')}")]

    if name == "move_file":
        file_id = arguments["file_id"]
        new_parent = arguments["new_parent_id"]
        f = service.files().get(fileId=file_id, fields="parents,name", supportsAllDrives=True).execute()
        old_parents = ",".join(f.get("parents", []))
        updated = service.files().update(fileId=file_id, addParents=new_parent, removeParents=old_parents, fields="id,name,parents", supportsAllDrives=True).execute()
        return [types.TextContent(type="text", text=f"Moved **{updated['name']}** to folder `{new_parent}`")]

    if name == "rename_file":
        file_id = arguments["file_id"]
        new_name = arguments["new_name"]
        updated = service.files().update(fileId=file_id, body={"name": new_name}, fields="id,name", supportsAllDrives=True).execute()
        return [types.TextContent(type="text", text=f"Renamed to **{updated['name']}** (ID: `{updated['id']}`)")]

    if name == "copy_file":
        file_id = arguments["file_id"]
        body = {}
        if arguments.get("new_name"): body["name"] = arguments["new_name"]
        if arguments.get("parent_id"): body["parents"] = [arguments["parent_id"]]
        copied = service.files().copy(fileId=file_id, body=body, fields="id,name,webViewLink", supportsAllDrives=True).execute()
        return [types.TextContent(type="text", text=f"Copied as **{copied['name']}**\n- ID: `{copied['id']}`\n- Link: {copied.get('webViewLink', '')}")]

    if name == "delete_file":
        file_id = arguments["file_id"]
        f = service.files().get(fileId=file_id, fields="name", supportsAllDrives=True).execute()
        service.files().update(fileId=file_id, body={"trashed": True}, supportsAllDrives=True).execute()
        return [types.TextContent(type="text", text=f"Moved **{f['name']}** to trash.")]

    if name == "share_file":
        file_id = arguments["file_id"]
        email = arguments["email"]
        role = arguments.get("role", "reader")
        permission = {"type": "user", "role": role, "emailAddress": email}
        service.permissions().create(fileId=file_id, body=permission, sendNotificationEmail=True).execute()
        f = service.files().get(fileId=file_id, fields="name", supportsAllDrives=True).execute()
        return [types.TextContent(type="text", text=f"Shared **{f['name']}** with {email} as **{role}**")]

    return [types.TextContent(type="text", text=f"Unknown tool: {name}")]


# =============================================================================
# Calendar dispatch
# =============================================================================

def _cal_dispatch(name: str, a: dict) -> list:
    cal = get_calendar_service()
    T = types.TextContent

    if name == "cal_list_calendars":
        result = cal.calendarList().list().execute()
        cals = result.get("items", [])
        out = [{"id": c["id"], "summary": c.get("summary", ""), "primary": c.get("primary", False), "accessRole": c.get("accessRole", ""), "backgroundColor": c.get("backgroundColor", "")} for c in cals]
        return [T(type="text", text=json.dumps(out, ensure_ascii=False, indent=2))]

    elif name == "cal_list_events":
        cal_id = a.get("calendar_id", "primary")
        now = datetime.now(timezone.utc)
        time_min = a.get("time_min", now.isoformat())
        time_max = a.get("time_max", (now + timedelta(days=30)).isoformat())
        max_results = min(a.get("max_results", 50), 2500)
        kwargs = {"calendarId": cal_id, "timeMin": time_min, "timeMax": time_max, "maxResults": max_results, "singleEvents": True, "orderBy": "startTime"}
        if a.get("query"): kwargs["q"] = a["query"]
        if a.get("show_deleted"): kwargs["showDeleted"] = True
        result = cal.events().list(**kwargs).execute()
        events = [_fmt_event(e) for e in result.get("items", [])]
        return [T(type="text", text=json.dumps(events, ensure_ascii=False, indent=2))]

    elif name == "cal_get_event":
        cal_id = a.get("calendar_id", "primary")
        event = cal.events().get(calendarId=cal_id, eventId=a["event_id"]).execute()
        return [T(type="text", text=json.dumps(_fmt_event(event), ensure_ascii=False, indent=2))]

    elif name == "cal_create_event":
        cal_id = a.get("calendar_id", "primary")
        body = _build_event_body(a)
        send = a.get("send_updates", "none")
        kwargs = {"calendarId": cal_id, "body": body, "sendUpdates": send}
        if "conferenceData" in body: kwargs["conferenceDataVersion"] = 1
        event = cal.events().insert(**kwargs).execute()
        return [T(type="text", text=json.dumps({"id": event["id"], "summary": event.get("summary"), "start": event.get("start"), "end": event.get("end"), "htmlLink": event.get("htmlLink")}, ensure_ascii=False, indent=2))]

    elif name == "cal_update_event":
        cal_id = a.get("calendar_id", "primary")
        event_id = a["event_id"]
        existing = cal.events().get(calendarId=cal_id, eventId=event_id).execute()
        updates = _build_event_body(a)
        existing.update(updates)
        send = a.get("send_updates", "none")
        kwargs = {"calendarId": cal_id, "eventId": event_id, "body": existing, "sendUpdates": send}
        if "conferenceData" in updates: kwargs["conferenceDataVersion"] = 1
        event = cal.events().update(**kwargs).execute()
        return [T(type="text", text=json.dumps({"id": event["id"], "summary": event.get("summary"), "start": event.get("start"), "end": event.get("end"), "htmlLink": event.get("htmlLink")}, ensure_ascii=False, indent=2))]

    elif name == "cal_delete_event":
        cal_id = a.get("calendar_id", "primary")
        send = a.get("send_updates", "none")
        cal.events().delete(calendarId=cal_id, eventId=a["event_id"], sendUpdates=send).execute()
        return [T(type="text", text=json.dumps({"ok": True, "deleted": a["event_id"]}))]

    elif name == "cal_move_event":
        cal_id = a.get("calendar_id", "primary")
        event = cal.events().move(calendarId=cal_id, eventId=a["event_id"], destination=a["destination_calendar_id"]).execute()
        return [T(type="text", text=json.dumps({"id": event["id"], "moved_to": a["destination_calendar_id"]}))]

    elif name == "cal_quick_add":
        cal_id = a.get("calendar_id", "primary")
        event = cal.events().quickAdd(calendarId=cal_id, text=a["text"]).execute()
        return [T(type="text", text=json.dumps({"id": event["id"], "summary": event.get("summary"), "start": event.get("start"), "end": event.get("end"), "htmlLink": event.get("htmlLink")}, ensure_ascii=False, indent=2))]

    elif name == "cal_freebusy":
        cal_ids = [c.strip() for c in a.get("calendar_ids", "primary").split(",") if c.strip()]
        body = {"timeMin": a["time_min"], "timeMax": a["time_max"], "items": [{"id": cid} for cid in cal_ids]}
        result = cal.freebusy().query(body=body).execute()
        calendars = result.get("calendars", {})
        out = {}
        for cid, data in calendars.items():
            busy = data.get("busy", [])
            out[cid] = {"busy_slots": len(busy), "busy": busy}
        return [T(type="text", text=json.dumps(out, ensure_ascii=False, indent=2))]

    elif name == "cal_get_colors":
        colors = cal.colors().get().execute()
        return [T(type="text", text=json.dumps(colors, ensure_ascii=False, indent=2))]

    elif name == "cal_list_recurring_instances":
        cal_id = a.get("calendar_id", "primary")
        kwargs = {"calendarId": cal_id, "eventId": a["event_id"]}
        if a.get("time_min"): kwargs["timeMin"] = a["time_min"]
        if a.get("time_max"): kwargs["timeMax"] = a["time_max"]
        if a.get("max_results"): kwargs["maxResults"] = a["max_results"]
        result = cal.events().instances(**kwargs).execute()
        events = [_fmt_event(e) for e in result.get("items", [])]
        return [T(type="text", text=json.dumps(events, ensure_ascii=False, indent=2))]

    elif name == "cal_respond_to_event":
        cal_id = a.get("calendar_id", "primary")
        event = cal.events().get(calendarId=cal_id, eventId=a["event_id"]).execute()
        me = cal.calendarList().get(calendarId=cal_id).execute().get("id", cal_id)
        for attendee in event.get("attendees", []):
            if attendee["email"] == me:
                attendee["responseStatus"] = a["response"]
                break
        else:
            event.setdefault("attendees", []).append({"email": me, "responseStatus": a["response"]})
        updated = cal.events().update(calendarId=cal_id, eventId=a["event_id"], body=event, sendUpdates="none").execute()
        return [T(type="text", text=json.dumps({"id": updated["id"], "response": a["response"]}))]

    return [T(type="text", text=f"Unknown calendar tool: {name}")]


# =============================================================================
# Sheets dispatch
# =============================================================================

_SHEETS_TOOLS = {
    "create_spreadsheet", "copy_spreadsheet", "get_spreadsheet_info", "list_spreadsheets",
    "share_spreadsheet", "get_permissions", "list_sheets", "add_sheet", "delete_sheet",
    "rename_sheet", "copy_sheet", "reorder_sheets", "read_range", "batch_read", "get_cell",
    "search_values", "get_sheet_as_table", "write_range", "append_rows", "batch_write",
    "update_cell", "clear_range", "insert_rows_columns", "delete_rows_columns", "format_cells",
    "merge_cells", "unmerge_cells", "set_column_width", "set_row_height", "auto_resize",
    "set_borders", "add_conditional_format", "set_data_validation", "create_named_range",
    "add_filter_view", "set_basic_filter", "create_pivot_table", "protect_range", "set_note",
    "get_notes", "export_spreadsheet", "get_formulas", "duplicate_as_template", "sort_range",
}


def _sheets_dispatch(name: str, a: dict) -> list:
    # File operations
    if name == "create_spreadsheet":
        result = _sheets_create_spreadsheet(title=a["title"], sheet_names=a.get("sheet_names"), folder_id=a.get("folder_id"))
    elif name == "copy_spreadsheet":
        result = _sheets_copy_spreadsheet(spreadsheet_id=a["spreadsheet_id"], new_title=a["new_title"], folder_id=a.get("folder_id"))
    elif name == "get_spreadsheet_info":
        result = _sheets_get_spreadsheet_info(a["spreadsheet_id"])
    elif name == "list_spreadsheets":
        result = _sheets_list_spreadsheets(query=a.get("query"), folder_id=a.get("folder_id"), max_results=a.get("max_results", 20), include_shared_drives=a.get("include_shared_drives", True))
    elif name == "share_spreadsheet":
        result = _sheets_share_spreadsheet(spreadsheet_id=a["spreadsheet_id"], email=a["email"], role=a["role"], send_notification=a.get("send_notification", True), message=a.get("message"))
    elif name == "get_permissions":
        result = _sheets_get_permissions(a["spreadsheet_id"])
    # Sheet management
    elif name == "list_sheets":
        result = _sheets_list_sheets(a["spreadsheet_id"])
    elif name == "add_sheet":
        result = _sheets_add_sheet(a["spreadsheet_id"], a["title"], rows=a.get("rows", 1000), columns=a.get("columns", 26), tab_color=a.get("tab_color"))
    elif name == "delete_sheet":
        result = _sheets_delete_sheet(a["spreadsheet_id"], a["sheet_id"])
    elif name == "rename_sheet":
        result = _sheets_rename_sheet(a["spreadsheet_id"], a["sheet_id"], a["new_title"])
    elif name == "copy_sheet":
        result = _sheets_copy_sheet(a["spreadsheet_id"], a["sheet_id"], a.get("destination_spreadsheet_id"))
    elif name == "reorder_sheets":
        result = _sheets_reorder_sheets(a["spreadsheet_id"], a["sheet_order"])
    # Reading
    elif name == "read_range":
        result = _sheets_read_range(a["spreadsheet_id"], a["range"], value_render=a.get("value_render", "FORMATTED_VALUE"))
    elif name == "batch_read":
        result = _sheets_batch_read(a["spreadsheet_id"], a["ranges"], value_render=a.get("value_render", "FORMATTED_VALUE"))
    elif name == "get_cell":
        result = _sheets_get_cell(a["spreadsheet_id"], a["cell"])
    elif name == "search_values":
        result = _sheets_search_values(a["spreadsheet_id"], a["query"], sheet_name=a.get("sheet_name"))
    elif name == "get_sheet_as_table":
        result = _sheets_get_sheet_as_table(a["spreadsheet_id"], sheet_name=a.get("sheet_name"), has_header=a.get("has_header", True), max_rows=a.get("max_rows", 0))
    # Writing
    elif name == "write_range":
        result = _sheets_write_range(a["spreadsheet_id"], a["range"], a["values"], value_input=a.get("value_input", "USER_ENTERED"))
    elif name == "append_rows":
        result = _sheets_append_rows(a["spreadsheet_id"], a["range"], a["values"], value_input=a.get("value_input", "USER_ENTERED"))
    elif name == "batch_write":
        result = _sheets_batch_write(a["spreadsheet_id"], a["data"], value_input=a.get("value_input", "USER_ENTERED"))
    elif name == "update_cell":
        result = _sheets_update_cell(a["spreadsheet_id"], a["cell"], a["value"])
    elif name == "clear_range":
        result = _sheets_clear_range(a["spreadsheet_id"], a["range"])
    elif name == "insert_rows_columns":
        result = _sheets_insert_rows_columns(a["spreadsheet_id"], a["sheet_id"], a["dimension"], a["start_index"], a["end_index"], inherit_from_before=a.get("inherit_from_before", False))
    elif name == "delete_rows_columns":
        result = _sheets_delete_rows_columns(a["spreadsheet_id"], a["sheet_id"], a["dimension"], a["start_index"], a["end_index"])
    # Formatting
    elif name == "format_cells":
        fmt_args = {k: v for k, v in a.items() if k not in ("spreadsheet_id", "sheet_id", "range")}
        result = _sheets_format_cells(a["spreadsheet_id"], a["sheet_id"], a["range"], **fmt_args)
    elif name == "merge_cells":
        result = _sheets_merge_cells(a["spreadsheet_id"], a["sheet_id"], a["range"], merge_type=a.get("merge_type", "MERGE_ALL"))
    elif name == "unmerge_cells":
        result = _sheets_unmerge_cells(a["spreadsheet_id"], a["sheet_id"], a["range"])
    elif name == "set_column_width":
        result = _sheets_set_column_width(a["spreadsheet_id"], a["sheet_id"], a["start_col"], a["end_col"], a["width"])
    elif name == "set_row_height":
        result = _sheets_set_row_height(a["spreadsheet_id"], a["sheet_id"], a["start_row"], a["end_row"], a["height"])
    elif name == "auto_resize":
        result = _sheets_auto_resize(a["spreadsheet_id"], a["sheet_id"], a["dimension"], start_index=a.get("start_index", 0), end_index=a.get("end_index", 26))
    elif name == "set_borders":
        result = _sheets_set_borders(a["spreadsheet_id"], a["sheet_id"], a["range"], style=a.get("style", "SOLID"), color=a.get("color", "#000000"), top=a.get("top", True), bottom=a.get("bottom", True), left=a.get("left", True), right=a.get("right", True), inner_horizontal=a.get("inner_horizontal", False), inner_vertical=a.get("inner_vertical", False))
    # Advanced
    elif name == "add_conditional_format":
        result = _sheets_add_conditional_format(a["spreadsheet_id"], a["sheet_id"], a["range"], a["rule_type"], values=a.get("values"), bg_color=a.get("bg_color"), text_color=a.get("text_color"), bold=a.get("bold"), min_color=a.get("min_color"), mid_color=a.get("mid_color"), max_color=a.get("max_color"))
    elif name == "set_data_validation":
        result = _sheets_set_data_validation(a["spreadsheet_id"], a["sheet_id"], a["range"], a["validation_type"], values=a.get("values"), strict=a.get("strict", True), input_message=a.get("input_message"))
    elif name == "create_named_range":
        result = _sheets_create_named_range(a["spreadsheet_id"], a["sheet_id"], a["name"], a["range"])
    elif name == "add_filter_view":
        result = _sheets_add_filter_view(a["spreadsheet_id"], a["sheet_id"], a["title"], a["range"], sort_column=a.get("sort_column"), sort_order=a.get("sort_order", "ASCENDING"), filter_column=a.get("filter_column"), hidden_values=a.get("hidden_values"))
    elif name == "set_basic_filter":
        if "range" in a and a["range"]:
            result = _sheets_set_basic_filter(a["spreadsheet_id"], a["sheet_id"], a["range"])
        else:
            result = _sheets_clear_basic_filter(a["spreadsheet_id"], a["sheet_id"])
    elif name == "create_pivot_table":
        result = _sheets_create_pivot_table(a["spreadsheet_id"], a["source_sheet_id"], a["source_range"], a["target_sheet_id"], target_row=a.get("target_row", 0), target_col=a.get("target_col", 0), row_fields=a["row_fields"], value_fields=a["value_fields"], column_fields=a.get("column_fields"), filter_fields=a.get("filter_fields"))
    # Protection & notes
    elif name == "protect_range":
        result = _sheets_protect_range(a["spreadsheet_id"], a["sheet_id"], range_a1=a.get("range"), description=a.get("description", ""), warning_only=a.get("warning_only", False), editors=a.get("editors"))
    elif name == "set_note":
        result = _sheets_set_note(a["spreadsheet_id"], a["sheet_id"], a["cell"], a["note"])
    elif name == "get_notes":
        result = _sheets_get_notes(a["spreadsheet_id"], a["sheet_name"])
    # Export & utilities
    elif name == "export_spreadsheet":
        result = _sheets_export_spreadsheet(a["spreadsheet_id"], export_format=a["format"], sheet_id=a.get("sheet_id"), download_dir=a.get("download_dir"))
    elif name == "get_formulas":
        result = _sheets_get_formulas(a["spreadsheet_id"], a["range"])
    elif name == "duplicate_as_template":
        result = _sheets_duplicate_as_template(a["spreadsheet_id"], a["new_title"], clear_data=a.get("clear_data", True), folder_id=a.get("folder_id"))
    elif name == "sort_range":
        result = _sheets_sort_range(a["spreadsheet_id"], a["sheet_id"], a["range"], a["sort_specs"])
    else:
        return _err(f"Unknown sheets tool: {name}")

    return _ok(result)


# =============================================================================
# Entry point
# =============================================================================

async def main():
    async with stdio_server() as (read_stream, write_stream):
        await server.run(read_stream, write_stream, server.create_initialization_options())


if __name__ == "__main__":
    asyncio.run(main())
